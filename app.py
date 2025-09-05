from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import psycopg2
import bcrypt
from io import BytesIO
from docx import Document
import fitz  # PyMuPDF
from fpdf import FPDF
from agent import analyze_document
from agent import get_answer_from_gemini, analyze_uploaded_document_content  # ✅ import your functions
from googletrans import Translator
#from agent import call_gemini_api 
from pydub import AudioSegment
import speech_recognition as sr
from datetime import datetime
import tempfile
from agent import analyze_document
from io import BytesIO
from docx import Document
import fitz  # PyMuPDF
from PIL import Image
import pytesseract
import os
from flask import Flask
import webbrowser
from threading import Timer
from flask import Flask, render_template
import smtplib
from email.mime.text import MIMEText
from datetime import datetime
from flask import Flask
import psycopg2
from datetime import date
import smtplib
from email.message import EmailMessage


import uuid
from flask import Flask, request, jsonify
from fpdf import FPDF
#from agent import generate_document_drafts
#from agent import model 
from flask import send_from_directory
import traceback


import re
import urllib.request
import urllib.error
import json as pyjson


import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
from dotenv import load_dotenv

# Create uploads directory if it doesn't exist
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
CORS(app)

# DB_CONFIG = {
#     "host": "localhost",
#     "port": 5432,
#     "dbname": "Document Drafting and Review Support Agent",
#     "user": "postgres",
#     "password": "Akif@Scaleable"
# }

# DB_CONFIG = {
#     "host": "legalassistantserver.postgres.database.azure.com",
#     "port": 5432,
#     "dbname": "Document_Drafting_and_Review_Support_Agent",  # Use the synced DB name
#     "user": "postgres@legalassistantserver",
#     "password": "Akif@Scaleable",
#     "sslmode": "require"  # Required for Azure PostgreSQL
# }

DB_CONFIG = {
    "host": "legalassistantserver.postgres.database.azure.com",
    "port": 5432,
    "dbname": "Document_Drafting_and_Review_Support_Agent",
    "user": "postgres@legalassistantserver",
    "password": "Akif@Scaleable",
    "sslmode": "require"
}





def get_connection():
    return psycopg2.connect(**DB_CONFIG)
def get_admin_email():
    """Get the first available admin email from the database"""
    try:
        conn = get_connection()
        cur = conn.cursor()
        cur.execute("SELECT email FROM admins LIMIT 1")
        admin_row = cur.fetchone()
        cur.close()
        conn.close()
        if admin_row:
            return admin_row[0]
    except Exception as e:
        print(f"⚠️ Could not fetch admin email from database: {e}")
    
    # Fallback to environment variable
    return os.getenv("ADMIN_EMAIL", "akif@scaleablesolutions.com")


def ensure_cases_customers_tables():
    """Ensure the cases and customers tables exist"""
    try:
        conn = get_connection()
        cur = conn.cursor()
        
        # Check if tables exist first
        cur.execute("""
            SELECT EXISTS (
                SELECT FROM information_schema.tables 
                WHERE table_name = 'cases'
            );
        """)
        cases_exists = cur.fetchone()[0]
        
        cur.execute("""
            SELECT EXISTS (
                SELECT FROM information_schema.tables 
                WHERE table_name = 'customer'
            );
        """)
        customer_exists = cur.fetchone()[0]
        
        print(f"🔍 Table check - Cases exists: {cases_exists}, Customer exists: {customer_exists}")
        
        # Only create tables if they don't exist
        if not customer_exists:
            # Create Customer table first (since Cases references it)
            cur.execute("""
                CREATE TABLE customer (
                    customerid SERIAL PRIMARY KEY,
                    name VARCHAR(255) NOT NULL,
                    contactinfo TEXT,
                    phone VARCHAR(20),
                    email VARCHAR(255) UNIQUE,
                    address TEXT,
                    profilefile VARCHAR(500),

                    -- Extra fields
                    dateofbirth DATE,
                    gender VARCHAR(20) CHECK (gender IN ('Male', 'Female', 'Other')),
                    cnic VARCHAR(20) UNIQUE,
                    occupation VARCHAR(100),
                    companyname VARCHAR(255),
                    emergencycontactname VARCHAR(255),
                    emergencycontactphone VARCHAR(20),
                    notes TEXT,

                    createddate TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    updateddate TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    document_type VARCHAR(20) CHECK (document_type IN ('Case', 'Customer')) DEFAULT 'Customer'
                );
            """)
            print("✅ Customer table created")
        else:
            # Add missing columns if table exists
            try:
                cur.execute("ALTER TABLE customer ADD COLUMN IF NOT EXISTS contactinfo TEXT;")
                cur.execute("ALTER TABLE customer ADD COLUMN IF NOT EXISTS address TEXT;")
                cur.execute("ALTER TABLE customer ADD COLUMN IF NOT EXISTS profilefile VARCHAR(500);")
                cur.execute("ALTER TABLE customer ADD COLUMN IF NOT EXISTS dateofbirth DATE;")
                cur.execute("ALTER TABLE customer ADD COLUMN IF NOT EXISTS gender VARCHAR(20);")
                cur.execute("ALTER TABLE customer ADD COLUMN IF NOT EXISTS emergencycontactname VARCHAR(255);")
                cur.execute("ALTER TABLE customer ADD COLUMN IF NOT EXISTS emergencycontactphone VARCHAR(20);")
                cur.execute("ALTER TABLE customer ADD COLUMN IF NOT EXISTS notes TEXT;")
                cur.execute("ALTER TABLE customer ADD COLUMN IF NOT EXISTS updateddate TIMESTAMP DEFAULT CURRENT_TIMESTAMP;")
                print("✅ Customer table columns updated")
                # Widen common fields to reduce 'value too long for type character varying(20)' errors
                try:
                    cur.execute("ALTER TABLE customer ALTER COLUMN phone TYPE VARCHAR(40);")
                except Exception:
                    pass
                try:
                    cur.execute("ALTER TABLE customer ALTER COLUMN emergencycontactphone TYPE VARCHAR(40);")
                except Exception:
                    pass
                try:
                    cur.execute("ALTER TABLE customer ALTER COLUMN cnic TYPE VARCHAR(40);")
                except Exception:
                    pass
            except Exception as e:
                print(f"ℹ️ Customer table columns update: {e}")
        
        # Migrate data from existing customers table if it exists
        try:
            cur.execute("""
                INSERT INTO customer (name, phone, email, cnic, occupation, companyname, createddate, document_type)
                SELECT name, phone, email, cnic, occupation, company, created_at, 
                       COALESCE(document_type, 'Customer')
                FROM customers
                WHERE NOT EXISTS (SELECT 1 FROM customer WHERE customer.name = customers.name)
            """)
            print("✅ Migrated existing customers data")
        except Exception as e:
            print(f"ℹ️ No existing customers to migrate or migration failed: {e}")
        
        if not cases_exists:
            # Create Cases table with the new schema provided by the user
            cur.execute("""
                CREATE TABLE cases (
                    caseid SERIAL PRIMARY KEY,
                    casename VARCHAR(255) NOT NULL,
                    description TEXT,
                    casetype VARCHAR(100),
                    casestatus VARCHAR(50) CHECK (casestatus IN ('Open', 'In Progress', 'Closed', 'On Hold')),
                    casepriority VARCHAR(20) CHECK (casepriority IN ('Low', 'Medium', 'High', 'Urgent')),
                    uploadfile VARCHAR(500),

                    -- Dates
                    createddate TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    openeddate DATE,
                    closeddate DATE,
                    duedate DATE,

                    -- Legal / Process
                    courtname VARCHAR(255),
                    hearingdate DATE,
                    judgename VARCHAR(255),
                    casecategory VARCHAR(100),
                    jurisdiction VARCHAR(255),

                    -- Financial
                    casefee NUMERIC(12,2),
                    billingstatus VARCHAR(20) CHECK (billingstatus IN ('Pending', 'Paid', 'Overdue')),

                    -- Relationship
                    customerid INT NOT NULL REFERENCES customer(customerid) ON DELETE CASCADE,
                    
                    -- Document Type
                    document_type VARCHAR(20) CHECK (document_type IN ('Case', 'Customer')) DEFAULT 'Case',
                    
                    -- Document Content
                    document_preview TEXT
                );
            """)
            print("✅ Cases table created")
        else:
            # Add missing columns if table exists
            try:
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS description TEXT;")
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS casetype VARCHAR(100);")
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS casestatus VARCHAR(50);")
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS casepriority VARCHAR(20);")
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS openeddate DATE;")
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS closeddate DATE;")
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS duedate DATE;")
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS courtname VARCHAR(255);")
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS hearingdate DATE;")
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS judgename VARCHAR(255);")
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS casecategory VARCHAR(100);")
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS jurisdiction VARCHAR(255);")
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS casefee NUMERIC(12,2);")
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS billingstatus VARCHAR(20);")
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS document_preview TEXT;")
                print("✅ Cases table columns updated")
            except Exception as e:
                print(f"ℹ️ Cases table columns update: {e}")
        
        # Debug: Check what columns actually exist
        cur.execute("""
            SELECT column_name, data_type 
            FROM information_schema.columns 
            WHERE table_name = 'cases'
            ORDER BY ordinal_position;
        """)
        columns = cur.fetchall()
        print("✅ Cases table columns:")
        for col in columns:
            print(f"  - {col[0]} ({col[1]})")
        
        # Create Customers table (keeping the existing schema for now, as no new schema was provided)
        cur.execute("""
            CREATE TABLE IF NOT EXISTS customers (
                id SERIAL PRIMARY KEY,
                name VARCHAR(255) NOT NULL,
                phone VARCHAR(50),
                email VARCHAR(255),
                cnic VARCHAR(50),
                occupation VARCHAR(255),
                company VARCHAR(255),
                document_type VARCHAR(20) DEFAULT 'Customer',
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
        """)
        
        # Add document_type column to customers table if it doesn't exist
        try:
            cur.execute("""
                ALTER TABLE customers 
                ADD COLUMN IF NOT EXISTS document_type VARCHAR(20) DEFAULT 'Customer';
            """)
        except Exception:
            pass  # Column might already exist
        
        conn.commit()
        cur.close()
        conn.close()
        print("✅ Cases and Customers tables ensured")
    except Exception as e:
        print(f"⚠️ Could not ensure cases/customers tables: {e}")

def ensure_documents_table():
    """Ensure the Documents table exists with proper schema for file URLs"""
    try:
        conn = get_connection()
        cur = conn.cursor()
        
        # Check if Documents table exists
        cur.execute("""
            SELECT EXISTS (
                SELECT FROM information_schema.tables 
                WHERE table_name = 'Documents'
            );
        """)
        table_exists = cur.fetchone()[0]
        
        if not table_exists:
            # Create Documents table with new schema
            cur.execute("""
                CREATE TABLE Documents (
                    DocumentID SERIAL PRIMARY KEY,
                    FileName VARCHAR(255) NOT NULL, 
                    FileURL VARCHAR(500),
                    UploadDate TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    -- Document Type (Customer or Case)
                    DocumentType VARCHAR(20) CHECK (DocumentType IN ('Customer', 'Case')),
                    -- Text content extracted from documents
                    extracted_text TEXT,
                    -- File content (binary)
                    content BYTEA,
                    -- File type
                    filetype VARCHAR(50),
                    -- Relationships
                    CustomerID INT REFERENCES Customer(CustomerID) ON DELETE CASCADE,
                    CaseID INT REFERENCES Cases(CaseID) ON DELETE CASCADE
                );
            """)
            print("✅ Documents table created with new schema")
        else:
            # Add missing columns if table exists
            try:
                cur.execute("ALTER TABLE Documents ADD COLUMN IF NOT EXISTS FileURL VARCHAR(500);")
                cur.execute("ALTER TABLE Documents ADD COLUMN IF NOT EXISTS DocumentType VARCHAR(20);")
                cur.execute("ALTER TABLE Documents ADD COLUMN IF NOT EXISTS extracted_text TEXT;")
                cur.execute("ALTER TABLE Documents ADD COLUMN IF NOT EXISTS content BYTEA;")
                cur.execute("ALTER TABLE Documents ADD COLUMN IF NOT EXISTS filetype VARCHAR(50);")
                cur.execute("ALTER TABLE Documents ADD COLUMN IF NOT EXISTS CustomerID INT REFERENCES Customer(CustomerID) ON DELETE CASCADE;")
                cur.execute("ALTER TABLE Documents ADD COLUMN IF NOT EXISTS CaseID INT REFERENCES Cases(CaseID) ON DELETE CASCADE;")
                
                # Add constraints if they don't exist
                try:
                    cur.execute("""
                        ALTER TABLE Documents 
                        ADD CONSTRAINT IF NOT EXISTS check_document_type 
                        CHECK (DocumentType IN ('Customer', 'Case'));
                    """)
                except Exception:
                    pass  # Constraint might already exist
                
                print("✅ Documents table columns updated")
            except Exception as e:
                print(f"ℹ️ Documents table columns update: {e}")
        
        conn.commit()
        cur.close()
        conn.close()
    except Exception as e:
        print(f"⚠️ Could not ensure documents table: {e}")

def save_file_and_get_url(file, folder_name):
    """Save uploaded file and return the URL"""
    try:
        # Generate unique filename
        filename = f"{uuid.uuid4()}_{file.filename}"
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], folder_name, filename)
        
        # Create folder if it doesn't exist
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        
        # Save file
        file.save(file_path)
        
        # Return URL (relative to the app)
        file_url = f"/uploads/{folder_name}/{filename}"
        return file_url
    except Exception as e:
        print(f"Error saving file: {e}")
        return None

def get_admin_emails():
    """Return all admin emails from the database; fall back to env ADMIN_EMAIL if table is empty."""
    try:
        ensure_admins_table()
        conn = get_connection(); cur = conn.cursor()
        cur.execute("SELECT email FROM admins WHERE email IS NOT NULL")
        rows = cur.fetchall()
        cur.close(); conn.close()
        emails = [row[0] for row in rows if row and row[0]]
        if emails:
            return emails
    except Exception as e:
        print(f"⚠️ Could not fetch admin emails from database: {e}")
    fallback = os.getenv("ADMIN_EMAIL")
    return [fallback] if fallback else []


# Optional settings table for dynamic Make.com webhooks
def ensure_settings_table():
    try:
        conn = get_connection(); cur = conn.cursor()
        cur.execute(
            """
            CREATE TABLE IF NOT EXISTS settings (
                key TEXT PRIMARY KEY,
                value TEXT
            );
            """
        )
        conn.commit(); cur.close(); conn.close()
    except Exception:
        pass

def get_setting(key: str) -> str:
    try:
        ensure_settings_table()
        conn = get_connection(); cur = conn.cursor()
        cur.execute("SELECT value FROM settings WHERE key=%s", (key,))
        row = cur.fetchone(); cur.close(); conn.close()
        if row and row[0]:
            return row[0]
    except Exception:
        pass
    return ""

def set_setting(key: str, value: str) -> None:
    try:
        ensure_settings_table()
        conn = get_connection(); cur = conn.cursor()
        cur.execute(
            "INSERT INTO settings (key, value) VALUES (%s, %s) ON CONFLICT (key) DO UPDATE SET value=EXCLUDED.value",
            (key, value),
        )
        conn.commit(); cur.close(); conn.close()
    except Exception:
        pass

@app.route("/test-make-config")
def test_make_config():
    admin_webhook = get_setting("MAKE_ADMIN_WEBHOOK") or os.getenv("MAKE_ADMIN_WEBHOOK") or os.getenv("MAKE_WEBHOOK_ADMIN", "")
    user_webhook = get_setting("MAKE_USER_WEBHOOK") or os.getenv("MAKE_USER_WEBHOOK") or os.getenv("MAKE_WEBHOOK_USER", "")
    return jsonify({"success": True, "admin_webhook": bool(admin_webhook), "user_webhook": bool(user_webhook)})

@app.route("/settings/webhook", methods=["POST"])
def set_webhook():
    data = request.get_json(silent=True) or {}
    key = (data.get("key") or "").strip()
    value = (data.get("value") or "").strip()
    if key not in {"MAKE_ADMIN_WEBHOOK", "MAKE_USER_WEBHOOK"}:
        return jsonify({"success": False, "message": "Invalid key"}), 400
    set_setting(key, value)
    return jsonify({"success": True})

# -------------------- Signup approval (admin-gated) --------------------
def ensure_signup_tables():
    try:
        conn = get_connection()
        cur = conn.cursor()
        # Pending signup requests table
        cur.execute(
            """
            CREATE TABLE IF NOT EXISTS signup_requests (
                id SERIAL PRIMARY KEY,
                username VARCHAR(100),
                email VARCHAR(150) UNIQUE,
                password_hash TEXT,
                status VARCHAR(20) DEFAULT 'pending',
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            """
        )
        # Ensure token column exists even if table was created before this feature
        cur.execute("ALTER TABLE signup_requests ADD COLUMN IF NOT EXISTS token TEXT UNIQUE;")
        conn.commit()
        cur.close()
        conn.close()
    except Exception as e:
        print("⚠️ Could not ensure signup tables:", e)

def backfill_missing_tokens():
    try:
        conn = get_connection()
        cur = conn.cursor()
        # Generate tokens for any rows that are pending and missing token
        cur.execute(
            """
            UPDATE signup_requests
            SET token = gen_random_uuid()::text
            WHERE token IS NULL AND status = 'pending';
            """
        )
        # If pgcrypto not available, fallback to Python-generated tokens per row
        if cur.rowcount == 0:
            cur.execute("SELECT id FROM signup_requests WHERE token IS NULL AND status='pending'")
            rows = cur.fetchall()
            for (rid,) in rows:
                cur.execute("UPDATE signup_requests SET token=%s WHERE id=%s", (str(uuid.uuid4()), rid))
        conn.commit()
        cur.close(); conn.close()
    except Exception as e:
        # Non-fatal
        print("ℹ️ Token backfill skipped:", e)

def _get_base_url() -> str:
    # Allow overriding base URL for links in emails (useful when running on a LAN or public host)
    return os.getenv("BASE_URL", "http://127.0.0.1:5050").rstrip("/")


def send_admin_approval_email(user_email: str, token: str | None, username: str | None = None):
    # Prefer token-based links; fall back to email-based if token is missing
    base_url = _get_base_url()
    if token:
        approve_link = f"{base_url}/approve?token={token}"
        reject_link = f"{base_url}/reject?token={token}"
    else:
        approve_link = f"{base_url}/approve?email={user_email}"
        reject_link = f"{base_url}/reject?email={user_email}"
    subject = f"New Signup Request: {username or ''} <{user_email}>".strip()
    body = (
        f"A new signup request has been submitted.\n\n"
        f"User: {username or ''}\n"
        f"Email: {user_email}\n\n"
        f"Approve: {approve_link}\n"
        f"Reject: {reject_link}\n"
    )
    html_body = f"""
    <div style='font-family:Arial,sans-serif;line-height:1.5'>
      <h2 style='margin:0 0 12px 0'>New Signup Request</h2>
      <p><strong>User:</strong> {username or ''}<br/>
         <strong>Email:</strong> {user_email}</p>
      <div style='margin-top:16px'>
        <a href='{approve_link}' style='background:#16a34a;color:#fff;padding:10px 14px;border-radius:6px;text-decoration:none;margin-right:8px;'>Approve</a>
        <a href='{reject_link}' style='background:#dc2626;color:#fff;padding:10px 14px;border-radius:6px;text-decoration:none;'>Reject</a>
      </div>
      <p style='color:#6b7280;margin-top:16px'>If the buttons do not work, use these links:<br/>
        <a href='{approve_link}'>{approve_link}</a><br/>
        <a href='{reject_link}'>{reject_link}</a></p>
    </div>
    """

    # Always print links for debugging
    print("\n===== Signup Approval Links =====")
    print("User:", username or "-")
    print("Email:", user_email)
    print("Approve:", approve_link)
    print("Reject:", reject_link)
    print("================================\n")
    try:
        # Try writing links to a local file as a fallback
        try:
            with open("approval_links.txt", "a", encoding="utf-8") as fp:
                fp.write(f"Email: {user_email}\nApprove: {approve_link}\nReject: {reject_link}\n---\n")
        except Exception:
            pass

        # Dynamically get admin emails from database
        admin_emails = get_admin_emails()
        print(f"📧 Using admin emails: {', '.join(admin_emails) if admin_emails else '(none)'}")

        # Prefer environment-configured SMTP credentials (support multiple env names)
        smtp_email = os.getenv("SENDER_EMAIL", "") or os.getenv("SMTP_EMAIL", "")
        smtp_pass = os.getenv("SENDER_PASSWORD", "") or os.getenv("SMTP_APP_PASSWORD", "")
        # Default to Gmail SMTP; override via env if needed
        smtp_server = os.getenv("SMTP_SERVER", "smtp.gmail.com")
        try:
            smtp_port = int(os.getenv("SMTP_PORT", "587"))
        except Exception:
            smtp_port = 587

        if smtp_email and smtp_pass:
            # Use SSL for port 465 (e.g., Gmail), otherwise STARTTLS for 587
            if smtp_port == 465:
                with smtplib.SMTP_SSL(smtp_server, smtp_port) as smtp:
                    smtp.login(smtp_email, smtp_pass)
                    for admin_email in admin_emails:
                        msg = EmailMessage()
                        msg["Subject"] = subject
                        # From handling: default to SMTP sender, optionally use user's email for display if enabled
                        email_from_mode = (os.getenv("EMAIL_FROM_MODE", "smtp") or "smtp").lower()
                        display_name = (username or user_email or smtp_email)
                        header_from = user_email if (email_from_mode == "user" and user_email) else smtp_email
                        msg["From"] = f"{display_name} <{header_from}>"
                        msg["To"] = admin_email
                        if user_email:
                            msg["Reply-To"] = user_email
                        msg.set_content(body)
                        msg.add_alternative(html_body, subtype="html")
                        smtp.send_message(msg)
            else:
                with smtplib.SMTP(smtp_server, smtp_port) as smtp:
                    smtp.starttls()
                    smtp.login(smtp_email, smtp_pass)
                    for admin_email in admin_emails:
                        msg = EmailMessage()
                        msg["Subject"] = subject
                        # From handling: default to SMTP sender, optionally use user's email for display if enabled
                        email_from_mode = (os.getenv("EMAIL_FROM_MODE", "smtp") or "smtp").lower()
                        display_name = (username or user_email or smtp_email)
                        header_from = user_email if (email_from_mode == "user" and user_email) else smtp_email
                        msg["From"] = f"{display_name} <{header_from}>"
                        msg["To"] = admin_email
                        if user_email:
                            msg["Reply-To"] = user_email
                        msg.set_content(body)
                        msg.add_alternative(html_body, subtype="html")
                        smtp.send_message(msg)
        else:
            # Try generic helper if configured
            try:
                for admin_email in admin_emails:
                    send_email(
                        to_email=admin_email,
                        subject=subject,
                        body=body,
                    )
            except Exception:
                pass
    except Exception as e:
        print("⚠️ Failed to send admin email:", e)

    # Also notify Make.com webhook if configured (admin notification)
    try:
        webhook = (
            get_setting("MAKE_ADMIN_WEBHOOK")
            or os.getenv("MAKE_ADMIN_WEBHOOK")
            or os.getenv("MAKE_WEBHOOK_ADMIN")
        )
        if webhook:
            payload = {
                "event": "signup_submitted",
                "username": username,
                "email": user_email,
                "token": token,
            }
            req = urllib.request.Request(webhook, data=pyjson.dumps(payload).encode("utf-8"), headers={"Content-Type": "application/json"})
            urllib.request.urlopen(req, timeout=5)
    except Exception as _:
        pass

@app.route("/signup", methods=["POST"])
def signup_request():
    # Accept JSON or form-encoded
    data = request.get_json(silent=True) or request.form
    username = (data.get("username") or "").strip()
    email = (data.get("email") or data.get("Email") or "").strip()
    password = (data.get("password") or "").strip()

    if not username or not email or not password:
        return jsonify({"status": "error", "message": "Missing required fields."}), 400

    try:
        ensure_signup_tables()
        token = str(uuid.uuid4())
        hashed_pw = bcrypt.hashpw(password.encode("utf-8"), bcrypt.gensalt()).decode("utf-8")

        conn = get_connection()
        cur = conn.cursor()
        try:
            # Upsert by email so the same email can re-request if previously rejected/pending
            cur.execute(
                """
                INSERT INTO signup_requests (username, email, password_hash, status, token)
                VALUES (%s, %s, %s, 'pending', %s)
                ON CONFLICT (email)
                DO UPDATE SET username = EXCLUDED.username,
                              password_hash = EXCLUDED.password_hash,
                              status = 'pending',
                              token = EXCLUDED.token
                RETURNING token;
                """,
                (username, email, hashed_pw, token),
            )
            saved_token = cur.fetchone()[0]
        except Exception as e:
            # If token column is still missing for any reason, fall back to schema without token
            print("ℹ️ Falling back to email-only approval (no token column):", e)
            cur.execute(
                """
                INSERT INTO signup_requests (username, email, password_hash, status)
                VALUES (%s, %s, %s, 'pending')
                ON CONFLICT (email)
                DO UPDATE SET username = EXCLUDED.username,
                              password_hash = EXCLUDED.password_hash,
                              status = 'pending'
                RETURNING email;
                """,
                (username, email, hashed_pw),
            )
            saved_token = None
        conn.commit()
        cur.close()
        conn.close()

        send_admin_approval_email(email, saved_token, username=username)
        return jsonify({"status": "pending", "message": "Signup request submitted for admin approval."})

    except Exception as e:
        return jsonify({"status": "error", "message": str(e)}), 500

@app.route("/approve", methods=["GET"])
def approve_request():
    token = (request.args.get("token") or "").strip()
    email_param = (request.args.get("email") or "").strip()
    if not token and not email_param:
        return "Invalid request", 400

    try:
        conn = get_connection()
        cur = conn.cursor()
        if token:
            cur.execute("SELECT username, email, password_hash, status FROM signup_requests WHERE token=%s", (token,))
        else:
            cur.execute("SELECT username, email, password_hash, status FROM signup_requests WHERE email=%s", (email_param,))
        row = cur.fetchone()
        if not row:
            cur.close(); conn.close()
            return "Invalid or expired token", 404

        username, email, password_hash, status = row
        if status == 'approved':
            cur.close(); conn.close()
            return "Already approved."

        # Insert into existing users table with email in the correct column
        try:
            cur.execute(
                "INSERT INTO users (username, email, password) VALUES (%s, %s, %s)",
                (username, email, password_hash),
            )
        except Exception:
            # Fallback for legacy schema where email column may not exist
            cur.execute(
                "INSERT INTO users (username, phone, password) VALUES (%s, %s, %s)",
                (username, email, password_hash),
            )
        if token:
            cur.execute("UPDATE signup_requests SET status='approved' WHERE token=%s", (token,))
        else:
            cur.execute("UPDATE signup_requests SET status='approved' WHERE email=%s", (email_param,))
        conn.commit()
        cur.close(); conn.close()
        # Send user notification email on approval
        try:
            send_user_status_email(email, username, "approved")
        except Exception as _e:
            # Non-fatal: log and continue
            print(f"⚠️ Failed to send user approval email: {_e}")
        # Notify Make.com webhook for user approval
        try:
            webhook = (
                get_setting("MAKE_USER_WEBHOOK")
                or os.getenv("MAKE_USER_WEBHOOK")
                or os.getenv("MAKE_WEBHOOK_USER")
            )
            if webhook:
                payload = {"event": "signup_approved", "username": username, "email": email}
                req = urllib.request.Request(webhook, data=pyjson.dumps(payload).encode("utf-8"), headers={"Content-Type": "application/json"})
                urllib.request.urlopen(req, timeout=5)
        except Exception:
            pass

        return "✅ Admin approved successfully."
    except Exception as e:
        return f"Error approving request: {str(e)}", 500

@app.route("/reject", methods=["GET"])
def reject_request():
    token = (request.args.get("token") or "").strip()
    email_param = (request.args.get("email") or "").strip()
    if not token and not email_param:
        return "Invalid request", 400
    try:
        conn = get_connection()
        cur = conn.cursor()
        # Fetch user info for email notification
        username = None
        email = None
        if token:
            cur.execute("SELECT username, email FROM signup_requests WHERE token=%s", (token,))
        else:
            cur.execute("SELECT username, email FROM signup_requests WHERE email=%s", (email_param,))
        row = cur.fetchone()
        if row:
            username, email = row
        # Perform rejection update
        if token:
            cur.execute("UPDATE signup_requests SET status='rejected' WHERE token=%s", (token,))
        else:
            cur.execute("UPDATE signup_requests SET status='rejected' WHERE email=%s", (email_param,))
        conn.commit()
        cur.close(); conn.close()
        # Send user notification email on rejection
        try:
            if email:
                send_user_status_email(email, username or email, "rejected")
        except Exception as _e:
            # Non-fatal: log and continue
            print(f"⚠️ Failed to send user rejection email: {_e}")
        # Notify Make.com webhook for rejection
        try:
            # Fetch email for context if token provided
            email_for_notify = email_param
            if token and not email_for_notify:
                try:
                    conn = get_connection(); cur = conn.cursor()
                    cur.execute("SELECT email FROM signup_requests WHERE token=%s", (token,))
                    r = cur.fetchone(); cur.close(); conn.close()
                    if r: email_for_notify = r[0]
                except Exception:
                    pass
            webhook = (
                get_setting("MAKE_USER_WEBHOOK")
                or os.getenv("MAKE_USER_WEBHOOK")
                or os.getenv("MAKE_WEBHOOK_USER")
            )
            if webhook and email_for_notify:
                payload = {"event": "signup_rejected", "email": email_for_notify}
                req = urllib.request.Request(webhook, data=pyjson.dumps(payload).encode("utf-8"), headers={"Content-Type": "application/json"})
                urllib.request.urlopen(req, timeout=5)
        except Exception:
            pass

        return "❌ User rejected."
    except Exception as e:
        return f"Error rejecting request: {str(e)}", 500

@app.route("/check-status", methods=["POST"])
def check_status():
    # Accept JSON or form
    data = request.get_json(silent=True) or request.form
    email = (data.get("email") or data.get("Email") or "").strip()
    if not email:
        return jsonify({"status": "error", "message": "Missing email"}), 400
    try:
        conn = get_connection()
        cur = conn.cursor()
        cur.execute("SELECT status FROM signup_requests WHERE email=%s", (email,))
        row = cur.fetchone()
        cur.close(); conn.close()
        if not row:
            return jsonify({"status": "not_found"})
        return jsonify({"status": row[0]})
    except Exception as e:
        return jsonify({"status": "error", "message": str(e)}), 500

@app.route("/admin/requests", methods=["GET"])
def admin_requests():
    # Simple HTML list of pending signup requests with Approve/Reject links
    try:
        ensure_signup_tables()
        backfill_missing_tokens()
        conn = get_connection()
        cur = conn.cursor()
        cur.execute(
            "SELECT id, username, email, status, COALESCE(token, '') AS token, created_at FROM signup_requests ORDER BY created_at DESC"
        )
        rows = cur.fetchall()
        cur.close(); conn.close()

        def row_html(r):
            _id, username, email, status, token, created_at = r
            if token:
                approve_link = f"/approve?token={token}"
                reject_link = f"/reject?token={token}"
            else:
                approve_link = f"/approve?email={email}"
                reject_link = f"/reject?email={email}"
            return f"""
            <tr>
              <td>{_id}</td>
              <td>{username or ''}</td>
              <td>{email}</td>
              <td>{status}</td>
              <td>{created_at}</td>
              <td>
                <a href='{approve_link}'>Approve</a> | <a href='{reject_link}'>Reject</a>
              </td>
            </tr>
            """

        rows_html = "\n".join(row_html(r) for r in rows)
        html = f"""
        <html>
          <head>
            <title>Signup Requests</title>
            <style>
              body {{ font-family: Arial, sans-serif; padding: 16px; }}
              table {{ border-collapse: collapse; width: 100%; }}
              th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
              th {{ background: #f5f5f5; }}
            </style>
          </head>
          <body>
            <h2>Signup Requests</h2>
            <table>
              <thead>
                <tr><th>ID</th><th>Username</th><th>Email</th><th>Status</th><th>Submitted</th><th>Actions</th></tr>
              </thead>
              <tbody>
                {rows_html}
              </tbody>
            </table>
          </body>
        </html>
        """
        return html
    except Exception as e:
        return f"Error loading admin page: {str(e)}", 500


@app.route("/admin/users", methods=["GET"])
def admin_users():
    try:
        ensure_signup_tables()
        conn = get_connection(); cur = conn.cursor()
        cur.execute("SELECT * FROM signup_requests WHERE status = 'approved' ORDER BY id DESC")
        rows = cur.fetchall()
        colnames = [desc[0] for desc in cur.description]
        cur.close(); conn.close()
        # Build table header and body dynamically
        thead = "".join([f"<th>{c}</th>" for c in colnames])
        body_rows = []
        for row in rows:
            tds = "".join([f"<td>{(str(val) if val is not None else '')}</td>" for val in row])
            body_rows.append(f"<tr>{tds}</tr>")
        rows_html = "\n".join(body_rows)
        return f"""
        <html>
          <head>
            <title>Approved Users</title>
            <style>
              body {{ font-family: Arial, sans-serif; padding: 16px; }}
              table {{ border-collapse: collapse; width: 100%; }}
              th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
              th {{ background: #f5f5f5; }}
            </style>
          </head>
          <body>
            <h2>Approved Users</h2>
            <table>
              <thead>
                <tr>{thead}</tr>
              </thead>
              <tbody>
                {rows_html}
              </tbody>
            </table>
          </body>
        </html>
        """
    except Exception as e:
        return f"Error loading users: {str(e)}", 500

@app.route("/admin/users.json", methods=["GET"])
def admin_users_json():
    try:
        ensure_signup_tables()
        conn = get_connection(); cur = conn.cursor()
        cur.execute("SELECT * FROM signup_requests WHERE status = 'approved' ORDER BY id DESC")
        rows = cur.fetchall()
        colnames = [desc[0] for desc in cur.description]
        cur.close(); conn.close()
        out = []
        for row in rows:
            obj = {}
            for idx, col in enumerate(colnames):
                val = row[idx]
                obj[col] = str(val) if isinstance(val, (datetime,)) else val
            out.append(obj)
        return jsonify({"users": out})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/admin/requests.json", methods=["GET"])
def admin_requests_json():
    try:
        ensure_signup_tables()
        conn = get_connection(); cur = conn.cursor()
        cur.execute("SELECT * FROM signup_requests ORDER BY created_at DESC")
        rows = cur.fetchall()
        colnames = [desc[0] for desc in cur.description]
        cur.close(); conn.close()
        out = []
        for row in rows:
            obj = {}
            for idx, col in enumerate(colnames):
                val = row[idx]
                obj[col] = str(val) if isinstance(val, (datetime,)) else val
            out.append(obj)
        return jsonify({"requests": out, "columns": colnames})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# --- Simple Admin Login (env-configured) ---
def ensure_admins_table():
    try:
        conn = get_connection(); cur = conn.cursor()
        cur.execute(
            """
            CREATE TABLE IF NOT EXISTS admins (
                id SERIAL PRIMARY KEY,
                username VARCHAR(100) NOT NULL,
                email VARCHAR(150) UNIQUE NOT NULL,
                password_hash TEXT NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            """
        )
        conn.commit(); cur.close(); conn.close()
    except Exception:
        # non-fatal
        pass

@app.route("/admin-login", methods=["POST"])
def admin_login():
    data = request.get_json(silent=True) or {}
    email = (data.get("email") or "").strip()
    password = (data.get("password") or "").strip()
    if not email or not password:
        return jsonify({"success": False, "message": "Missing email or password"}), 400
    try:
        ensure_admins_table()
        conn = get_connection(); cur = conn.cursor()
        cur.execute("SELECT username, password_hash FROM admins WHERE email=%s", (email,))
        row = cur.fetchone(); cur.close(); conn.close()
        if not row:
            return jsonify({"success": False, "message": "Admin not found"}), 404
        stored_username, password_hash = row
        if bcrypt.checkpw(password.encode("utf-8"), (password_hash or '').encode("utf-8")):
            return jsonify({"success": True, "username": stored_username, "email": email})
        return jsonify({"success": False, "message": "Invalid admin credentials"}), 401
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

@app.route("/admin-signup", methods=["POST"])
def admin_signup():
    data = request.get_json(silent=True) or {}
    username = (data.get("username") or "").strip()
    email = (data.get("email") or "").strip()
    password = (data.get("password") or "").strip()
    if not username or not email or not password:
        return jsonify({"success": False, "message": "Missing required fields."}), 400
    try:
        ensure_admins_table()
        hashed = bcrypt.hashpw(password.encode("utf-8"), bcrypt.gensalt()).decode("utf-8")
        conn = get_connection(); cur = conn.cursor()
        cur.execute("INSERT INTO admins (username, email, password_hash) VALUES (%s, %s, %s)", (username, email, hashed))
        conn.commit(); cur.close(); conn.close()
        return jsonify({"success": True, "message": "Admin registered. You can sign in now."})
    except psycopg2.errors.UniqueViolation:
        return jsonify({"success": False, "message": "Admin with this email already exists."}), 409
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

@app.route("/admin/test-email", methods=["GET"])
def admin_test_email():
    try:
        # Resolve admin recipients dynamically
        admin_emails = get_admin_emails()
        if not admin_emails:
            return jsonify({"success": False, "message": "No admin emails found in database or ADMIN_EMAIL env."}), 400

        # SMTP configuration (support multiple env names)
        smtp_email = os.getenv("SENDER_EMAIL", "") or os.getenv("SMTP_EMAIL", "")
        smtp_pass = os.getenv("SENDER_PASSWORD", "") or os.getenv("SMTP_APP_PASSWORD", "")
        smtp_server = os.getenv("SMTP_SERVER", "smtp.gmail.com")
        try:
            smtp_port = int(os.getenv("SMTP_PORT", "587"))
        except Exception:
            smtp_port = 587

        if not smtp_email or not smtp_pass:
            return jsonify({"success": False, "message": "SMTP not configured. Set SMTP_EMAIL and SMTP_APP_PASSWORD (and optionally SMTP_SERVER/SMTP_PORT)."}), 400

        subject = "Test Email: Admin Notification System"
        body = "This is a test email to confirm SMTP and the dynamic admin list are working."
        html_body = "<div style='font-family:Arial,sans-serif'><h3>Admin Email Test</h3><p>If you received this, SMTP is configured and admin list resolution works.</p></div>"

        # Send to all admins
        if smtp_port == 465:
            with smtplib.SMTP_SSL(smtp_server, smtp_port) as smtp:
                smtp.login(smtp_email, smtp_pass)
                for admin_email in admin_emails:
                    msg = EmailMessage()
                    msg["Subject"] = subject
                    msg["From"] = smtp_email
                    msg["To"] = admin_email
                    msg.set_content(body)
                    msg.add_alternative(html_body, subtype="html")
                    smtp.send_message(msg)
        else:
            with smtplib.SMTP(smtp_server, smtp_port) as smtp:
                smtp.starttls()
                smtp.login(smtp_email, smtp_pass)
                for admin_email in admin_emails:
                    msg = EmailMessage()
                    msg["Subject"] = subject
                    msg["From"] = smtp_email
                    msg["To"] = admin_email
                    msg.set_content(body)
                    msg.add_alternative(html_body, subtype="html")
                    smtp.send_message(msg)

        print(f"📧 Test email sent to admins: {', '.join(admin_emails)} via {smtp_server}:{smtp_port}")
        return jsonify({"success": True, "admins": admin_emails, "smtp_server": smtp_server, "smtp_port": smtp_port})
    except Exception as e:
        print("⚠️ Failed to send test email:", e)
        return jsonify({"success": False, "message": str(e)}), 500



@app.route("/admin/approve", methods=["POST"])
def admin_approve_user():
    data = request.get_json(silent=True) or {}
    request_id = data.get("request_id")
    email = data.get("email")
    
    if not request_id or not email:
        return jsonify({"success": False, "message": "Missing request_id or email"}), 400
    
    try:
        conn = get_connection()
        cur = conn.cursor()
        
        # Get user details from signup_requests
        cur.execute("SELECT username, password_hash FROM signup_requests WHERE id=%s AND email=%s", (request_id, email))
        row = cur.fetchone()
        if not row:
            cur.close(); conn.close()
            return jsonify({"success": False, "message": "Request not found"}), 404
        
        username, password_hash = row
        
        # Check if already approved
        cur.execute("SELECT status FROM signup_requests WHERE id=%s", (request_id,))
        status_row = cur.fetchone()
        if status_row and status_row[0] == 'approved':
            cur.close(); conn.close()
            return jsonify({"success": False, "message": "Request already approved"}), 400
        
        # Insert into users table: prefer `email` column, fallback to legacy `phone`
        try:
            cur.execute(
                "INSERT INTO users (username, email, password) VALUES (%s, %s, %s)",
                (username, email, password_hash),
            )
        except Exception:
            cur.execute(
                "INSERT INTO users (username, phone, password) VALUES (%s, %s, %s)",
                (username, email, password_hash),
            )
        
        # Update status to approved
        cur.execute("UPDATE signup_requests SET status='approved' WHERE id=%s", (request_id,))

        conn.commit()
        cur.close()
        conn.close()
         # Send email to user after approval
        send_user_status_email(email, username, "approved")
        # Send approval email to user via Make.com
        try:
            webhook = (
                get_setting("MAKE_USER_WEBHOOK")
                or os.getenv("MAKE_USER_WEBHOOK")
                or os.getenv("MAKE_WEBHOOK_USER")
            )
            if webhook:
                payload = {"event": "signup_approved", "username": username, "email": email}
                req = urllib.request.Request(webhook, data=pyjson.dumps(payload).encode("utf-8"), headers={"Content-Type": "application/json"})
                urllib.request.urlopen(req, timeout=5)
        except Exception:
            pass
        
        return jsonify({"success": True, "message": "Admin approved successfully"})
        
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

@app.route("/admin/reject", methods=["POST"])
def admin_reject_user():
    data = request.get_json(silent=True) or {}
    request_id = data.get("request_id")
    email = data.get("email")
    
    if not request_id or not email:
        return jsonify({"success": False, "message": "Missing request_id or email"}), 400
    
    try:
        conn = get_connection()
        cur = conn.cursor()
        
        # Update status to rejected
        cur.execute("UPDATE signup_requests SET status='rejected' WHERE id=%s AND email=%s", (request_id, email))
        
        if cur.rowcount == 0:
            cur.close(); conn.close()
            return jsonify({"success": False, "message": "Request not found"}), 404
        

  
        conn.commit()
        cur.close()
        conn.close()
         # Send email to user after rejection
        # --- inside admin_reject_user after commit ---
        send_user_status_email(email, email.split("@")[0], "rejected")

        try:
            webhook = os.getenv("MAKE_USER_WEBHOOK") or os.getenv("MAKE_WEBHOOK_USER")
            if webhook:
                payload = {"event": "signup_rejected", "email": email}
                req = urllib.request.Request(webhook, data=pyjson.dumps(payload).encode("utf-8"), headers={"Content-Type": "application/json"})
                urllib.request.urlopen(req, timeout=5)
        except Exception:
            pass
        
        return jsonify({"success": False, "message": "User rejected successfully"})
        
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500




# Reg
@app.route("/register", methods=["POST"])
def register():
    data = request.get_json()
    username = data.get("username")
    email = data.get("email") or data.get("Email")
    password = data.get("password")

    if not username or not password:
        return jsonify({"success": False, "message": "Username and password are required"}), 400

    hashed_password = bcrypt.hashpw(password.encode("utf-8"), bcrypt.gensalt()).decode("utf-8")

    try:
        conn = get_connection()
        cur = conn.cursor()
        cur.execute(
            "INSERT INTO users (username, email, password) VALUES (%s, %s, %s)",
            (username, email, hashed_password)
        )
        conn.commit()
        cur.close()
        conn.close()
        return jsonify({"success": True, "message": "Registration successful."})
    except psycopg2.errors.UniqueViolation:
        return jsonify({"success": False, "message": "Username already exists."}), 409
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

@app.route("/login", methods=["POST"])
def login():
    data = request.get_json()
    # Expect email for login only (no username fallback)
    email = (data.get("email") or data.get("Email") or "").strip()
    password = (data.get("password") or "").strip()

    if not email or not password:
        return jsonify({"success": False, "message": "Email and password are required"}), 400

    try:
        conn = get_connection()
        cur = conn.cursor()
        # Check in signup_requests table for approved users
        cur.execute(
            "SELECT password_hash FROM signup_requests WHERE LOWER(email) = LOWER(%s) AND status = 'approved' ORDER BY id DESC LIMIT 1",
            (email,),
        )
        result = cur.fetchone()
        cur.close(); conn.close()

        if not result:
            return jsonify({"success": False, "message": "User not found or not approved."}), 404

        stored_hash = (result[0] or "").strip()
        if not stored_hash:
            return jsonify({"success": False, "message": "Invalid account."}), 401

        try:
            valid = bcrypt.checkpw(password.encode("utf-8"), stored_hash.encode("utf-8"))
        except Exception:
            # If stored hash is already bytes
            try:
                valid = bcrypt.checkpw(password.encode("utf-8"), stored_hash)
            except Exception:
                valid = False
        if valid:
            return jsonify({"success": True, "message": "Login successful."})
        else:
            return jsonify({"success": False, "message": "Invalid password."}), 401

    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500




# --- robust_gmail_email.py ---
import os
import smtplib
import ssl
from email.message import EmailMessage
from dotenv import load_dotenv

# Load environment variables from .env
load_dotenv()

def send_admin_approval_email(user_email, token=None, username=None):
    """
    Sends an approval request email to all admins when a new user signs up.
    Requires EMAIL_ADDRESS and EMAIL_PASSWORD in .env (Gmail App Password).
    """

    # ✅ Prepare approval/reject links
    approve_link = f"http://127.0.0.1:5050/approve?token={token}" if token else "N/A"
    reject_link = f"http://127.0.0.1:5050/reject?token={token}" if token else "N/A"

    # ✅ Email subject and body
    subject = "New Signup Request – Approval Needed"
    body = f"""
A new user has requested to sign up:

Username: {username}
Email: {user_email}

Approve: {approve_link}
Reject: {reject_link}
"""
    html_body = f"""
<h2>New Signup Request</h2>
<p><b>Username:</b> {username}</p>
<p><b>Email:</b> {user_email}</p>
<p>
    <a href="{approve_link}">✅ Approve</a> |
    <a href="{reject_link}">❌ Reject</a>
</p>
"""

    # ✅ Fetch admin emails from your database
    try:
        conn = get_connection()  # Replace with your DB connection function
        cur = conn.cursor()
        cur.execute("SELECT email FROM admins")
        admins = [row[0] for row in cur.fetchall()]
        cur.close()
        conn.close()
    except Exception as e:
        print(f"⚠️ Failed to fetch admin emails: {e}")
        admins = []

    if not admins:
        print("⚠️ No admins found in DB. Email skipped.")
        return

    # ✅ Load Gmail credentials from .env
    email_address = os.getenv("EMAIL_ADDRESS")
    email_password = os.getenv("EMAIL_PASSWORD")

    if not email_address or not email_password:
        print("❌ EMAIL_ADDRESS or EMAIL_PASSWORD not set in .env")
        return

    # ✅ Send email to each admin
    for admin_email in admins:
        msg = EmailMessage()
        msg["From"] = email_address
        msg["To"] = admin_email
        msg["Subject"] = subject
        msg.set_content(body)
        msg.add_alternative(html_body, subtype="html")

        try:
            # Using Gmail SMTP SSL with port 465
            with smtplib.SMTP_SSL("smtp.gmail.com", 465, context=ssl.create_default_context()) as smtp:
                smtp.login(email_address, email_password)
                smtp.send_message(msg)
            print(f"✅ Email sent to admin: {admin_email}")
        except Exception as e:
            print(f"❌ Email sending failed for {admin_email}: {e}")




# User email setup
def send_user_status_email(user_email, username, status):
    """
    Send status email to user after admin approves or rejects.
    status = 'approved' | 'rejected'
    """
    # Use same SMTP config as admin emails
    smtp_email = os.getenv("SENDER_EMAIL", "") or os.getenv("SMTP_EMAIL", "") or os.getenv("EMAIL_ADDRESS", "")
    smtp_pass = os.getenv("SENDER_PASSWORD", "") or os.getenv("SMTP_APP_PASSWORD", "") or os.getenv("EMAIL_PASSWORD", "")
    smtp_server = os.getenv("SMTP_SERVER", "smtp.gmail.com")
    try:
        smtp_port = int(os.getenv("SMTP_PORT", "587"))
    except Exception:
        smtp_port = 587

    if not smtp_email or not smtp_pass:
        print("❌ SMTP sender credentials not set (SENDER_EMAIL/SMTP_EMAIL + SENDER_PASSWORD/SMTP_APP_PASSWORD)")
        return

    subject = "Your Signup Request Status"
    if status == "approved":
        body = f"""
Hello {username},

✅ Congratulations! Your signup request has been approved.
You can now log in to your account using your credentials.

Regards,
Admin Team
"""
        html_body = f"""
<h2>Signup Approved 🎉</h2>
<p>Hello <b>{username}</b>,</p>
<p>✅ Congratulations! Your signup request has been <b>approved</b>.</p>
<p>You can now <a href="http://127.0.0.1:5050">Login here</a>.</p>
"""
    else:
        body = f"""
Hello {username},

❌ Sorry, your signup request has been rejected.
If you believe this is a mistake, please contact support.

Regards,
Admin Team
"""
        html_body = f"""
<h2>Signup Rejected ❌</h2>
<p>Hello <b>{username}</b>,</p>
<p>Unfortunately, your signup request has been <b>rejected</b>.</p>
<p>If you think this was an error, please contact support.</p>
"""

    msg = EmailMessage()
    msg["From"] = smtp_email
    # Ensure only the intended user receives this email
    msg["To"] = user_email
    if "Cc" in msg:
        del msg["Cc"]
    if "Bcc" in msg:
        del msg["Bcc"]
    msg["Subject"] = subject
    msg.set_content(body)
    msg.add_alternative(html_body, subtype="html")

    try:
        print(f"📧 Preparing to send user status email to: {user_email} (status={status}) via {smtp_server}:{smtp_port}")
        if smtp_port == 465:
            with smtplib.SMTP_SSL(smtp_server, smtp_port, context=ssl.create_default_context()) as smtp:
                smtp.login(smtp_email, smtp_pass)
                smtp.send_message(msg, from_addr=smtp_email, to_addrs=[user_email])
        else:
            with smtplib.SMTP(smtp_server, smtp_port) as smtp:
                smtp.starttls()
                smtp.login(smtp_email, smtp_pass)
                smtp.send_message(msg, from_addr=smtp_email, to_addrs=[user_email])
        print(f"✅ Status email sent to user: {user_email}")
    except Exception as e:
        print(f"❌ Failed to send status email to {user_email}: {e}")


# ================== FILE UPLOAD ================== #
@app.route("/upload", methods=["POST"])
def upload():
    if 'files[]' not in request.files:
        return jsonify({"error": "No files provided"}), 400

    uploaded_files = request.files.getlist('files[]')
    if not uploaded_files:
        return jsonify({"error": "Empty file list"}), 400

    responses = []
    for file in uploaded_files:
        if file.filename == '':
            responses.append({"error": "Empty filename"})
            continue

        try:
            filename = file.filename
            ext = os.path.splitext(filename)[1].lower()
            filetype = ext[1:]
            file_bytes = file.read()

            # Extract text
            if ext == ".txt":
                extracted_text = file_bytes.decode("utf-8")
            elif ext == ".docx":
                doc = Document(BytesIO(file_bytes))
                extracted_text = "\n".join([para.text for para in doc.paragraphs])
            elif ext == ".pdf":
                pdf = fitz.open(stream=file_bytes, filetype="pdf")
                extracted_text = "\n".join([page.get_text() for page in pdf])
            elif ext in [".jpg", ".jpeg", ".png"]:
                image = Image.open(BytesIO(file_bytes))
                extracted_text = pytesseract.image_to_string(image)
            elif ext in [".ppt", ".pptx"]:
                try:
                    # Try to import pptx, but don't fail if not available
                    try:
                        from pptx import Presentation  # type: ignore
                        prs = Presentation(BytesIO(file_bytes))
                        extracted_text = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    extracted_text.append(shape.text)
                        extracted_text = "\n".join(extracted_text)
                    except ImportError:
                        extracted_text = "PowerPoint file - text extraction requires python-pptx package. Install with: pip install python-pptx"
                except Exception as e:
                    extracted_text = f"PowerPoint file - extraction error: {str(e)}"
            else:
                responses.append({"filename": filename, "error": "Unsupported file type."})
                continue

            # Save to DB (be resilient to schema differences). Also try to persist a preview.
            doc_id = None
            try:
                conn = get_connection(); cur = conn.cursor()
                # Ensure optional columns exist
                try:
                    cur.execute("ALTER TABLE documents ADD COLUMN IF NOT EXISTS document_preview TEXT;")
                except Exception:
                    pass
                try:
                    # Preferred schema with doc_type and document_preview
                    cur.execute(
                        """
                        INSERT INTO documents (filename, filetype, content, extracted_text, document_preview, doc_type)
                        VALUES (%s, %s, %s, %s, %s, %s)
                        RETURNING id
                        """,
                        (filename, filetype, file_bytes, extracted_text, extracted_text, 'general')
                    )
                except Exception:
                    # Fallback legacy schema without doc_type
                    conn.rollback()
                    try:
                        # Try with document_preview only
                        cur.execute(
                            """
                            INSERT INTO documents (filename, filetype, content, extracted_text, document_preview)
                            VALUES (%s, %s, %s, %s, %s)
                            RETURNING id
                            """,
                            (filename, filetype, file_bytes, extracted_text, extracted_text)
                        )
                    except Exception:
                        # Final fallback minimal schema
                        conn.rollback()
                        cur.execute(
                            """
                            INSERT INTO documents (filename, filetype, content, extracted_text)
                            VALUES (%s, %s, %s, %s)
                            RETURNING id
                            """,
                            (filename, filetype, file_bytes, extracted_text)
                        )
                doc_id = cur.fetchone()[0]
                conn.commit()
                cur.close(); conn.close()
            except Exception as db_err:
                # Ensure connection closed on error if opened
                try:
                    cur.close(); conn.close()
                except Exception:
                    pass
                # Do not fail extraction because of DB issues; return text with error info
                responses.append({
                    "filename": filename,
                    "error": str(db_err),
                    "text": extracted_text
                })
                continue

            responses.append({"filename": filename, "id": doc_id, "text": extracted_text})

        except Exception as e:
            responses.append({"filename": file.filename, "error": str(e)})

    return jsonify({"results": responses})


# ================== DOCUMENT ANALYSIS ================== #
@app.route("/analyze", methods=["POST"])
def analyze():
    data = request.get_json()
    content = data.get("content")
    case_id = data.get("case_id")
    if not content:
        return jsonify({"error": "No content provided."}), 400
    try:
        analysis = analyze_document(content)
        # Persist extracted/analysis text if a case_id is provided
        if case_id:
            try:
                conn = get_connection(); cur = conn.cursor()
                cur.execute("ALTER TABLE cases ADD COLUMN IF NOT EXISTS extracted_data TEXT;")
                cur.execute("UPDATE cases SET extracted_data=%s WHERE caseid=%s", (content, case_id))
                conn.commit(); cur.close(); conn.close()
            except Exception as _e:
                print("⚠️ Failed to save extracted_data to cases:", _e)
        return jsonify({"analysis": analysis})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ================== PDF GENERATION ================== #
class SafePDF(FPDF):
    def __init__(self):
        super().__init__()
        self.set_auto_page_break(auto=True, margin=15)

def sanitize_text(text):
    return text.replace('\r', '').strip()

def break_long_words(text, max_length=100):
    return ' '.join([word if len(word) <= max_length else word[:max_length] + '...' for word in text.split()])

def generate_combined_pdf(original, analysis):
    pdf = SafePDF()
    pdf.add_page()
    page_width = pdf.w - 2 * pdf.l_margin

    def render_section(title, content):
        pdf.set_font("Helvetica", "B", 14)
        pdf.cell(0, 10, title, ln=True)
        pdf.set_font("Helvetica", "", 12)
        for line in content.strip().split('\n'):
            line = sanitize_text(line)
            line = break_long_words(line)
            pdf.multi_cell(page_width, 8, line)

    render_section("Original Document Content", original)
    pdf.add_page()
    render_section("AI Analysis", analysis)

    buffer = BytesIO()
    pdf.output(buffer)
    buffer.seek(0)
    return buffer, "report.pdf", "application/pdf"


@app.route("/download", methods=["POST"])
def download():
    data = request.get_json()
    original = data.get("original")
    analysis = data.get("analysis")
    format = data.get("format")

    if not all([original, analysis, format]):
        return jsonify({"error": "Missing data for download."}), 400

    if format == "txt":
        buffer = BytesIO(f"=== Original Document ===\n{original}\n\n=== AI Analysis ===\n{analysis}".encode("utf-8"))
        return send_file(buffer, as_attachment=True, download_name="report.txt", mimetype="text/plain")

    elif format == "docx":
        doc = Document()
        doc.add_heading("Original Document", level=1)
        for line in original.splitlines(): doc.add_paragraph(line)
        doc.add_page_break()
        doc.add_heading("AI Analysis", level=1)
        for line in analysis.splitlines(): doc.add_paragraph(line)
        buffer = BytesIO(); doc.save(buffer); buffer.seek(0)
        return send_file(buffer, as_attachment=True, download_name="report.docx", mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    
    elif format == "pdf":
        buffer, filename, mime = generate_combined_pdf(original, analysis)
        return send_file(buffer, as_attachment=True, download_name=filename, mimetype=mime)

    return jsonify({"error": "Unsupported file format."}), 400


# ================== CHAT ================== #
from deep_translator import GoogleTranslator
translator = GoogleTranslator(source="auto", target="en")
#from agent import get_answer_from_legislation

from flask import request, jsonify
from googletrans import Translator


translator = Translator()  # initialize globally

from flask import Flask, request, jsonify
#from agent import get_answer_from_llama, build_legislation_query
from googletrans import Translator

translator = Translator()

@app.route("/chat", methods=["POST"])
def chat():
    data = request.get_json()
    question = data.get("question", "")
    document = data.get("document", "")  # Get document preview from frontend

    if not question:
        return jsonify({"answer": "❌ Missing question."}), 400

    try:
        # Step 1: Build context - prioritize provided document, then fall back to database
        context_parts = []
        
        # Add provided document preview if available (highest priority)
        if document and document.strip():
            context_parts.append(f"UPLOADED DOCUMENT PREVIEW:\n{document}\n---")
        
        # Fetch comprehensive context from Cases table and Case documents
        conn = get_connection()
        cur = conn.cursor()
        
        # Make sure preview columns exist
        
        try:
            cur.execute("ALTER TABLE Documents ADD COLUMN IF NOT EXISTS document_preview TEXT;")
        except Exception:
            pass

        # Fetch all case data with customer information (include extracted_data and document_preview)
        cur.execute("""
            SELECT 
                c.CaseID,
                c.CaseName,
                c.Description,
                c.CaseType,
                c.CaseStatus,
                c.CasePriority,
                c.CourtName,
                c.HearingDate,
                c.JudgeName,
                c.CaseCategory,
                c.Jurisdiction,
                c.CaseFee,
                c.BillingStatus,
                c.CreatedDate,
                c.OpenedDate,
                c.ClosedDate,
                c.DueDate,
                cust.Name as CustomerName,
                cust.Phone as CustomerPhone,
                cust.Email as CustomerEmail,
                c.extracted_data,
                c.document_preview
            FROM Cases c
            LEFT JOIN Customer cust ON c.CustomerID = cust.CustomerID
            ORDER BY c.CreatedDate DESC
        """)
        case_rows = cur.fetchall()
        
        # Fetch case-related documents WITH extracted/document_preview text content
        # Include any rows that are linked to a CaseID and have some text, regardless of DocumentType value
        cur.execute("""
            SELECT 
                d.FileName,
                d.FileURL,
                d.UploadDate,
                d.DocumentType,
                d.CaseID,
                COALESCE(d.document_preview, d.extracted_text) as doc_preview
            FROM Documents d
            WHERE d.CaseID IS NOT NULL
              AND COALESCE(d.document_preview, d.extracted_text) IS NOT NULL
            ORDER BY d.UploadDate DESC
        """)
        case_doc_rows = cur.fetchall()
        
        # Fetch general documents WITH extracted/document_preview text content (not tied to a case)
        cur.execute("""
            SELECT FileName, FileURL, UploadDate, DocumentType, COALESCE(document_preview, extracted_text) as doc_preview
            FROM Documents 
            WHERE (CaseID IS NULL OR CaseID = 0)
              AND COALESCE(document_preview, extracted_text) IS NOT NULL
        """)
        general_doc_rows = cur.fetchall()
        
        cur.close()
        conn.close()
        
        # Add case information
        if case_rows:
            case_context = []
            for case in case_rows:
                # Prefer case-level document_preview; fall back to extracted_data
                case_text = (case[21] or case[20] or '')
                case_info = f"""
Case ID: {case[0]}
Case Name: {case[1] or 'N/A'}
Description: {case[2] or 'N/A'}
Case Type: {case[3] or 'N/A'}
Status: {case[4] or 'N/A'}
Priority: {case[5] or 'N/A'}
Court: {case[6] or 'N/A'}
Hearing Date: {case[7] or 'N/A'}
Judge: {case[8] or 'N/A'}
Category: {case[9] or 'N/A'}
Jurisdiction: {case[10] or 'N/A'}
Case Fee: {case[11] or 'N/A'}
Billing Status: {case[12] or 'N/A'}
Created: {case[13] or 'N/A'}
Opened: {case[14] or 'N/A'}
Closed: {case[15] or 'N/A'}
Due Date: {case[16] or 'N/A'}
Customer: {case[17] or 'N/A'}
Customer Phone: {case[18] or 'N/A'}
Customer Email: {case[19] or 'N/A'}
Case Document Preview:
{case_text}
---"""
                case_context.append(case_info)
            context_parts.append("CASES INFORMATION:\n" + "\n".join(case_context))
        
        # Add case documents WITH actual content
        if case_doc_rows:
            doc_context = []
            for doc in case_doc_rows:
                doc_info = f"""
Document: {doc[0]} (Case ID: {doc[4]}, Uploaded: {doc[2]})
Document Type: {doc[3]}
Document Content:
{doc[5] or 'No text content available'}
---"""
                doc_context.append(doc_info)
            context_parts.append("CASE DOCUMENTS WITH CONTENT:\n" + "\n".join(doc_context))
        
        # Add general documents WITH actual content
        if general_doc_rows:
            doc_context = []
            for doc in general_doc_rows:
                doc_info = f"""
Document: {doc[0]} (Type: {doc[3]}, Uploaded: {doc[2]})
Document Content:
{doc[4] or 'No text content available'}
---"""
                doc_context.append(doc_info)
            context_parts.append("GENERAL DOCUMENTS WITH CONTENT:\n" + "\n".join(doc_context))
        
        combined_text = "\n\n".join(context_parts)

        if not combined_text.strip():
            # Provide a helpful response when no data is available
            return jsonify({"answer": "👋 Hello! I'm your legal case assistant. Currently, there are no cases or documents in the database yet. You can:\n\n1. Add a new case using the '+ Add New Case' button\n2. Upload case documents\n3. Ask me general legal questions\n\nOnce you have cases and documents, I'll be able to provide detailed information about them!"})

        # Step 3: Translate question → English
        detected = translator.detect(question)
        source_lang = detected.lang
        question_en = translator.translate(question, src=source_lang, dest='en').text

        # Step 4: Get AI answer with enhanced context (prioritizing uploaded document)
        try:
            # If user provided specific document content, use document-focused analysis
            if document and document.strip() and len(document.strip()) > 100:
                print(f"📄 Using document-focused analysis for question: {question_en[:100]}...")
                answer_en = analyze_uploaded_document_content(document, question_en)
            else:
                # Use general context analysis
                print(f"📋 Using general context analysis for question: {question_en[:100]}...")
            answer_en = get_answer_from_gemini(question_en, combined_text)
        except Exception as e:
            print(f"❌ Error getting AI answer: {e}")
            # Fallback response when AI service fails
            if "hi" in question_en.lower() or "hello" in question_en.lower():
                answer_en = "👋 Hello! I'm your legal case assistant. I can help you with questions about your cases, documents, and legal matters. How can I assist you today?"
            else:
                answer_en = "🤖 I'm here to help with your legal cases and documents. Currently, I can see the information in your database. What would you like to know?"

        # Step 5: Translate back to user language
        try:
            final_answer = translator.translate(answer_en, src='en', dest=source_lang).text
        except Exception as e:
            print(f"❌ Translation error: {e}")
            final_answer = answer_en  # Use English if translation fails

        return jsonify({"answer": final_answer})

    except Exception as e:
        print(f"❌ Error during chat: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/chat-case", methods=["POST"])
def chat_case():
    """Enhanced chat endpoint specifically for case-related questions"""
    data = request.get_json()
    question = data.get("question", "")

    if not question:
        return jsonify({"answer": "❌ Missing question."}), 400

    try:
        # Step 1: Fetch comprehensive context from Cases table and Case documents
        conn = get_connection()
        cur = conn.cursor()
        
        # Fetch all case data with customer information
        cur.execute("""
            SELECT 
                c.CaseID,
                c.CaseName,
                c.Description,
                c.CaseType,
                c.CaseStatus,
                c.CasePriority,
                c.CourtName,
                c.HearingDate,
                c.JudgeName,
                c.CaseCategory,
                c.Jurisdiction,
                c.CaseFee,
                c.BillingStatus,
                c.CreatedDate,
                c.OpenedDate,
                c.ClosedDate,
                c.DueDate,
                cust.Name as CustomerName,
                cust.Phone as CustomerPhone,
                cust.Email as CustomerEmail
            FROM Cases c
            LEFT JOIN Customer cust ON c.CustomerID = cust.CustomerID
            ORDER BY c.CreatedDate DESC
        """)
        case_rows = cur.fetchall()
        
        # Fetch case-related documents WITH extracted text content
        cur.execute("""
            SELECT 
                d.FileName,
                d.FileURL,
                d.UploadDate,
                d.DocumentType,
                d.CaseID,
                d.extracted_text
            FROM Documents d
            WHERE d.DocumentType = 'Case' AND d.CaseID IS NOT NULL
            ORDER BY d.UploadDate DESC
        """)
        case_doc_rows = cur.fetchall()
        
        cur.close()
        conn.close()

        # Step 2: Build comprehensive context
        context_parts = []
        
        # Add case information
        if case_rows:
            case_context = []
            for case in case_rows:
                case_info = f"""
Case ID: {case[0]}
Case Name: {case[1] or 'N/A'}
Description: {case[2] or 'N/A'}
Case Type: {case[3] or 'N/A'}
Status: {case[4] or 'N/A'}
Priority: {case[5] or 'N/A'}
Court: {case[6] or 'N/A'}
Hearing Date: {case[7] or 'N/A'}
Judge: {case[8] or 'N/A'}
Category: {case[9] or 'N/A'}
Jurisdiction: {case[10] or 'N/A'}
Case Fee: {case[11] or 'N/A'}
Billing Status: {case[12] or 'N/A'}
Created: {case[13] or 'N/A'}
Opened: {case[14] or 'N/A'}
Closed: {case[15] or 'N/A'}
Due Date: {case[16] or 'N/A'}
Customer: {case[17] or 'N/A'}
Customer Phone: {case[18] or 'N/A'}
Customer Email: {case[19] or 'N/A'}
---"""
                case_context.append(case_info)
            context_parts.append("CASES INFORMATION:\n" + "\n".join(case_context))
        
        # Add case documents WITH actual content
        if case_doc_rows:
            doc_context = []
            for doc in case_doc_rows:
                doc_info = f"""
Document: {doc[0]} (Case ID: {doc[4]}, Uploaded: {doc[2]})
Document Type: {doc[3]}
Document Content:
{doc[5] or 'No text content available'}
---"""
                doc_context.append(doc_info)
            context_parts.append("CASE DOCUMENTS WITH CONTENT:\n" + "\n".join(doc_context))
        
        combined_text = "\n\n".join(context_parts)

        if not combined_text.strip():
            return jsonify({"answer": "🤖 I don't have any case information in my database yet. To help you better, please:\n\n1. Add some cases using the 'Case' button\n2. Add some customers using the 'Customer' button\n3. Upload relevant documents for your cases\n\nOnce you have some data, I'll be able to answer your questions about your cases and documents!"})

        # Step 3: Translate question → English
        detected = translator.detect(question)
        source_lang = detected.lang
        question_en = translator.translate(question, src=source_lang, dest='en').text

        # Step 4: Get AI answer with enhanced case context
        answer_en = get_answer_from_gemini(question_en, combined_text)

        # Step 5: Translate back to user language
        final_answer = translator.translate(answer_en, src='en', dest=source_lang).text

        return jsonify({"answer": final_answer})

    except Exception as e:
        print("❌ Error during case chat:", e)
        return jsonify({"error": str(e)}), 500


# voice upload and transcription 
@app.route("/upload-voice", methods=["POST"])
def upload_voice():
    try:
        if "audio" not in request.files:
            print("❌ No audio in request.files")
            return jsonify({"error": "No audio file in request"}), 400

        file = request.files["audio"]
        if file.filename == "":
            print("❌ Empty filename")
            return jsonify({"error": "No selected file"}), 400

        # Read audio content
        content = file.read()
        filename = file.filename
        print(f"✅ Received file: {filename}, size: {len(content)} bytes")

        # Save to PostgreSQL
        conn = get_connection()
        cur = conn.cursor()
        cur.execute(
            "INSERT INTO voice_files (filename, content) VALUES (%s, %s) RETURNING id",
            (filename, psycopg2.Binary(content))
        )
        voice_id = cur.fetchone()[0]
        conn.commit()
        cur.close()
        conn.close()

        print(f"✅ Voice file saved with ID {voice_id}")
        return jsonify({"voice_id": voice_id})

    except Exception as e:
        print("❌ Exception during voice upload:", str(e))
        return jsonify({"error": str(e)}), 500



@app.route("/voice-summary/<int:voice_id>", methods=["GET"])
def get_voice_summary(voice_id):
    try:
        # Fetch audio bytes from PostgreSQL
        conn = get_connection()
        cur = conn.cursor()
        cur.execute("SELECT content FROM voice_files WHERE id = %s", (voice_id,))
        result = cur.fetchone()
        cur.close()
        conn.close()

        if not result:
            return jsonify({"error": "Voice file not found."}), 404

        audio_bytes = result[0]

        # Save original audio to a temporary MP3 file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as temp_input:
            temp_input.write(audio_bytes)
            temp_input.flush()
            input_path = temp_input.name

        # Convert to WAV (mono, 16-bit, 16000Hz)
        sound = AudioSegment.from_file(input_path)
        sound = sound.set_channels(1)
        sound = sound.set_frame_rate(16000)

        # Create a temp WAV output file path
        temp_output = tempfile.NamedTemporaryFile(delete=False, suffix=".wav")
        output_path = temp_output.name
        temp_output.close()  # Close so it's not locked

        sound.export(output_path, format="wav")

        # Transcribe the audio using speech_recognition
        recognizer = sr.Recognizer()
        with sr.AudioFile(output_path) as source:
            audio_data = recognizer.record(source)
            try:
                text = recognizer.recognize_google(audio_data)
                print("✅ Transcription:", text)
            except sr.UnknownValueError:
                return jsonify({"error": "Google Speech Recognition could not understand the audio."}), 400
            except sr.RequestError as e:
                return jsonify({"error": f"Google Speech API error: {str(e)}"}), 500

        # Clean up temporary files
        os.remove(input_path)
        os.remove(output_path)

        # Generate summary from transcribed text
        summary = analyze_document(text)

        return jsonify({"summary": summary})

    except Exception as e:
        return jsonify({"error": f"Voice processing failed: {str(e)}"}), 500

# vocie part
@app.route("/ask-voice", methods=["POST"])
def ask_voice():
    data = request.get_json()
    voice_question = data.get("question", "")
    document = data.get("document", "")

    # Allow asking without a document; we'll fall back to DB context
    if not voice_question:
        return jsonify({"error": "Missing question"}), 400

    try:
        # Detect language and translate to English for processing
        detected = translator.detect(voice_question)
        source_lang = getattr(detected, 'lang', 'en') or 'en'
        question_en = translator.translate(voice_question, src=source_lang, dest='en').text if source_lang != 'en' else voice_question

        # Prepare context: prefer provided document; otherwise use DB context (documents + cases)
        context_en = ""
        if document:
            # Use provided document directly (already in English or will be processed)
                context_en = document
        else:
            # Build combined context from DB when no document is provided
            try:
                conn = get_connection()
                cur = conn.cursor()
                # Fetch documents with extracted text content
                cur.execute("""
                    SELECT FileName, FileURL, UploadDate, DocumentType, extracted_text 
                    FROM Documents 
                    WHERE DocumentType IS NOT NULL AND extracted_text IS NOT NULL
                """)
                doc_rows = cur.fetchall()
                try:
                    cur.execute("SELECT description FROM cases WHERE description IS NOT NULL")
                    case_rows = cur.fetchall()
                except Exception:
                    case_rows = []
                cur.close(); conn.close()
                parts = []
                if doc_rows:
                    doc_info = []
                    for r in doc_rows:
                        if r and r[0]:
                            doc_info.append(f"""
Document: {r[0]} (Type: {r[3]}, Uploaded: {r[2]})
Document Content:
{r[4] or 'No text content available'}
---""")
                    parts.append("\n\n".join(doc_info))
                if case_rows:
                    parts.append("\n\n".join([r[0] for r in case_rows if r and r[0]]))
                context_en = "\n\n".join([p for p in parts if p])
            except Exception:
                context_en = ""

        # Use document-focused analysis if document is provided, otherwise general analysis
        if document and document.strip() and len(document.strip()) > 100:
            print(f"📄 Voice: Using document-focused analysis for question: {question_en[:100]}...")
            answer_en = analyze_uploaded_document_content(document, question_en)
        else:
            print(f"📋 Voice: Using general context analysis for question: {question_en[:100]}...")
        answer_en = get_answer_from_gemini(question_en, context_en)
        final_answer = translator.translate(answer_en, src='en', dest=source_lang).text if source_lang != 'en' else answer_en
        return jsonify({"answer": final_answer, "lang": source_lang})
    except Exception as e:
        return jsonify({"error": f"Failed to process voice question: {str(e)}"}), 500


# Process audio file for case analysis
@app.route("/process-audio", methods=["POST"])
def process_audio():
    try:
        if "audio" not in request.files:
            print("❌ No audio in request.files")
            return jsonify({"success": False, "error": "No audio file in request"}), 400

        file = request.files["audio"]
        if file.filename == "":
            print("❌ Empty filename")
            return jsonify({"success": False, "error": "No selected file"}), 400

        # Read audio content
        content = file.read()
        filename = file.filename
        print(f"✅ Received audio file: {filename}, size: {len(content)} bytes")

        # Save original audio to a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as temp_input:
            temp_input.write(content)
            temp_input.flush()
            input_path = temp_input.name

        # Convert to WAV (mono, 16-bit, 16000Hz) for better speech recognition
        sound = AudioSegment.from_file(input_path)
        sound = sound.set_channels(1)
        sound = sound.set_frame_rate(16000)

        # Create a temp WAV output file path
        temp_output = tempfile.NamedTemporaryFile(delete=False, suffix=".wav")
        output_path = temp_output.name
        temp_output.close()  # Close so it's not locked

        sound.export(output_path, format="wav")

        # Transcribe the audio using speech_recognition
        recognizer = sr.Recognizer()
        with sr.AudioFile(output_path) as source:
            audio_data = recognizer.record(source)
            try:
                text = recognizer.recognize_google(audio_data)
                print("✅ Transcription:", text)
            except sr.UnknownValueError:
                return jsonify({"success": False, "error": "Google Speech Recognition could not understand the audio."}), 400
            except sr.RequestError as e:
                return jsonify({"success": False, "error": f"Google Speech API error: {str(e)}"}), 500

        # Clean up temporary files
        os.remove(input_path)
        os.remove(output_path)

        # Generate comprehensive analysis from transcribed text using the document analysis function
        if text and text.strip():
            print("📄 Generating analysis for transcribed text...")
            summary = analyze_document(text)
            print("✅ Analysis generated successfully")
            return jsonify({"success": True, "summary": summary})
        else:
            return jsonify({"success": False, "error": "No text was transcribed from the audio."}), 400

    except Exception as e:
        print(f"❌ Exception during audio processing: {str(e)}")
        return jsonify({"success": False, "error": f"Audio processing failed: {str(e)}"}), 500


# ✅ Tell pytesseract where Tesseract is installed
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"


# for reminder email
def send_email(to_email, subject, body):
    msg = EmailMessage()
    msg['Subject'] = subject
    msg['From'] = "your_email@example.com"
    msg['To'] = to_email
    msg.set_content(body)

    with smtplib.SMTP('smtp.gmail.com', 587) as smtp:
        smtp.starttls()
        smtp.login("your_email@example.com", "your_app_password")
        smtp.send_message(msg)

@app.route("/check-deadlines", methods=["GET"])
def check_deadlines():
    try:
        conn = get_connection()
        cur = conn.cursor()
        
        # Check for cases with upcoming deadlines instead
        cur.execute("""
            SELECT c.CaseName, c.DueDate, c.CaseStatus, c.CasePriority
            FROM Cases c
            WHERE c.DueDate < %s AND c.CaseStatus IN ('Open', 'In Progress')
            ORDER BY c.DueDate ASC
        """, (date.today(),))
        overdue_cases = cur.fetchall()
        cur.close()
        conn.close()

        if overdue_cases:
            message = "⚠️ The following cases have deadlines that have passed:\n\n"
            for case_name, due_date, status, priority in overdue_cases:
                message += f"• {case_name} (Due: {due_date}, Status: {status}, Priority: {priority})\n"
            
            return jsonify({"message": message, "overdue_count": len(overdue_cases)})
        else:
            return jsonify({"message": "✅ No overdue cases found.", "overdue_count": 0})
            
    except Exception as e:
        return jsonify({"error": str(e)}), 500



# Application & Document Drafting
import google.generativeai as genai

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Gemini API setup
genai.configure(api_key="AIzaSyBRAKfYEnbImVZOEESX7KuIA8Op5mWI9js")  # Replace with your Gemini API key
model = genai.GenerativeModel("gemini-2.5-flash-lite")  # or gemini-1.5-pro if preferred

# Templates
app.template_folder = "templates"
app.static_folder = "uploads"



# Generate documents from PDF
@app.route("/generate-drafts", methods=["POST"])
def generate_drafts():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "Empty filename"}), 400

    try:
        file_path = os.path.join(UPLOAD_FOLDER, file.filename)
        file.save(file_path)

        doc = fitz.open(file_path)
        text = "".join([page.get_text() for page in doc])

        if not text.strip():
            return jsonify({"error": "No readable text found."}), 400

        # Prompt Gemini
        prompt1 = f"Generate a professional cover letter based on this document:\n{text[:3000]}"
        prompt2 = f"Extract important form-like data from this document:\n{text[:3000]}"
        prompt3 = f"Create a supporting statement based on this document:\n{text[:3000]}"

        response1 = model.generate_content(prompt1)
        response2 = model.generate_content(prompt2)
        response3 = model.generate_content(prompt3)

        return jsonify({
            "cover_letter": getattr(response1, 'text', '❌ No cover letter text.'),
            "form_data": getattr(response2, 'text', '❌ No form data text.'),
            "supporting": getattr(response3, 'text', '❌ No supporting statement text.')
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({
            "cover_letter": "❌ Error generating cover letter.",
            "form_data": "❌ Error extracting form data.",
            "supporting": "❌ Error generating supporting statement.",
            "debug": str(e)
        }), 500

# Generate downloadable PDF
@app.route("/generate-pdf", methods=["POST"])
def generate_pdf():
    data = request.get_json()
    title = data.get("title", "Document")
    content = data.get("content", "")

    if not content.strip():
        return jsonify({"error": "No content provided."}), 400

    try:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", "B", 14)
        pdf.cell(0, 10, title, ln=True)
        pdf.ln(5)
        pdf.set_font("Arial", size=12)

        for line in content.split("\n"):
            pdf.multi_cell(0, 10, line.strip())

        filename = f"{uuid.uuid4()}.pdf"
        path = os.path.join(UPLOAD_FOLDER, filename)
        pdf.output(path)

        return jsonify({"url": f"/uploads/{filename}"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# Serve generated PDF - Route removed to avoid conflict




# ------------------- ✨ Form Template Filling Functions (New) ✨ -------------------
from flask import Flask, render_template, request, send_file, jsonify
import os
import shutil
import re
from difflib import SequenceMatcher
import fitz
from PIL import Image
import docx
from pdfrw import PdfReader, PdfWriter, PdfDict
import google.generativeai as genai
import json
from typing import List, Dict  # <-- ADD THIS LINE

UPLOAD_DIR = "uploads"
FILLED_DIR = "filled_pdfs"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(FILLED_DIR, exist_ok=True)

# Reuse your existing variables and functions
extracted_data = {}
form_fields = []

# --- OCR & Text Extraction ---
def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        doc = fitz.open(pdf_path)
        for page in doc:
            text_content = page.get_text()
            if text_content.strip():
                text += text_content + "\n"
            else:
                ocr_text = try_ocr_extraction(page)
                if ocr_text:
                    text += ocr_text + "\n"
        doc.close()
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return ""
    return text.strip()

def extract_text_from_docx(docx_path):
    try:
        doc = docx.Document(docx_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return ""

def extract_text_from_file(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()
    if file_extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension == '.docx':
        return extract_text_from_docx(file_path)
    else:
        print(f"Unsupported file format: {file_extension}")
        return ""

def try_ocr_extraction(page):
    try:
        import pytesseract
        try:
            pytesseract.get_tesseract_version()
        except Exception:
            print("Tesseract OCR not available. Skipping OCR extraction.")
            return ""
        pix = page.get_pixmap()
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        ocr_text = pytesseract.image_to_string(img)
        return ocr_text.strip()
    except ImportError:
        print("pytesseract not available. Skipping OCR extraction.")
        return ""
    except Exception as e:
        print(f"OCR extraction failed: {e}")
        return ""

# --- Form Field Extraction ---
def extract_form_fields_from_template(template_text: str) -> List[str]:
    field_pattern = r'\{\{([^}]+)\}\}'
    matches = re.findall(field_pattern, template_text)
    fields = []
    for match in matches:
        field_name = match.strip().lower()
        field_name = clean_field_name(field_name)
        if field_name and field_name not in fields:
            fields.append(field_name)
    return fields

def clean_field_name(field_name: str) -> str:
    template_words = ['template', 'form', 'field', 'input', 'placeholder']
    words = field_name.split()
    cleaned_words = [word for word in words if word not in template_words]
    cleaned = ' '.join(cleaned_words).strip()
    cleaned = re.sub(r'[^\w\s]', '', cleaned)
    return cleaned

def extract_all_data_from_client_document(client_text: str) -> Dict[str, str]:
    if not client_text.strip():
        return {}
    extracted_data = {}
    extracted_data.update(extract_structured_key_value_pairs(client_text))
    extracted_data.update(extract_using_common_patterns(client_text))
    extracted_data.update(extract_from_document_sections(client_text))
    extracted_data.update(extract_standalone_values(client_text))
    return extracted_data

def extract_structured_key_value_pairs(text: str) -> Dict[str, str]:
    result = {}
    lines = text.split('\n')
    for line in lines:
        line = line.strip()
        if not line or ':' not in line:
            continue
        parts = line.split(':', 1)
        if len(parts) != 2:
            continue
        key = parts[0].strip().lower()
        value = parts[1].strip()
        key = clean_field_name(key)
        if key and value and is_valid_value(value):
            result[key] = value
    return result

def extract_using_common_patterns(text: str) -> Dict[str, str]:
    result = {}
    # Name patterns
    name_patterns = [
        r'name[:\s]*([a-zA-Z\s]+)',
        r'full name[:\s]*([a-zA-Z\s]+)',
        r'client name[:\s]*([a-zA-Z\s]+)',
        r'applicant[:\s]*([a-zA-Z\s]+)',
        r'^([a-zA-Z\s]+)$'
    ]
    for pattern in name_patterns:
        match = re.search(pattern, text, re.IGNORECASE | re.MULTILINE)
        if match:
            value = match.group(1).strip()
            if is_valid_value(value) and not any(char.isdigit() for char in value):
                result['name'] = value
                break
    # Date patterns
    date_patterns = [
        r'date of birth[:\s]*(\d{1,2}[/-]\d{1,2}[/-]\d{2,4})',
        r'dob[:\s]*(\d{1,2}[/-]\d{1,2}[/-]\d{2,4})',
        r'born[:\s]*(\d{1,2}[/-]\d{1,2}[/-]\d{2,4})',
        r'birth date[:\s]*(\d{1,2}[/-]\d{1,2}[/-]\d{2,4})',
        r'(\d{1,2}[/-]\d{1,2}[/-]\d{2,4})'
    ]
    for pattern in date_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            value = match.group(1).strip()
            if is_valid_value(value):
                result['dob'] = value
                break
    # Address patterns
    address_patterns = [
        r'address[:\s]*([^.\n]+)',
        r'residence[:\s]*([^.\n]+)',
        r'lives at[:\s]*([^.\n]+)',
        r'home address[:\s]*([^.\n]+)',
        r'street[:\s]*([^.\n]+)'
    ]
    for pattern in address_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            value = match.group(1).strip()
            if is_valid_value(value) and len(value) > 5:
                result['address'] = value
                break
    # Employment patterns
    employment_patterns = [
        r'employment[:\s]*([^.\n]+)',
        r'job[:\s]*([^.\n]+)',
        r'occupation[:\s]*([^.\n]+)',
        r'works at[:\s]*([^.\n]+)',
        r'position[:\s]*([^.\n]+)',
        r'title[:\s]*([^.\n]+)'
    ]
    for pattern in employment_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            value = match.group(1).strip()
            if is_valid_value(value) and len(value) > 2:
                result['employment'] = value
                break
    # Income patterns
    income_patterns = [
        r'income[:\s]*([^.\n]+)',
        r'salary[:\s]*([^.\n]+)',
        r'earnings[:\s]*([^.\n]+)',
        r'annual income[:\s]*([^.\n]+)',
        r'monthly income[:\s]*([^.\n]+)',
        r'(\$[\d,]+)',
        r'(£[\d,]+)',
        r'(\d+[\d,]*\s*(?:per annum|per year|monthly|weekly))'
    ]
    for pattern in income_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            value = match.group(1).strip()
            if is_valid_value(value) and len(value) > 1:
                result['income'] = value
                break
    return result

def extract_from_document_sections(text: str) -> Dict[str, str]:
    result = {}
    lines = text.split('\n')
    current_section = ""
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.isupper() or (len(line) < 50 and line.endswith(':')):
            current_section = line.lower().replace(':', '').strip()
            continue
        if current_section and ':' in line:
            parts = line.split(':', 1)
            if len(parts) == 2:
                key = parts[0].strip().lower()
                value = parts[1].strip()
                combined_key = f"{current_section}_{key}"
                if is_valid_value(value):
                    result[combined_key] = value
    return result

def extract_standalone_values(text: str) -> Dict[str, str]:
    result = {}
    email_pattern = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
    emails = re.findall(email_pattern, text)
    if emails:
        result['email'] = emails[0]
    phone_pattern = r'(\+?[\d\s\-\(\)]{10,})'
    phones = re.findall(phone_pattern, text)
    if phones:
        result['phone'] = phones[0].strip()
    postal_pattern = r'\b[A-Z]{1,2}\d[A-Z\d]?\s?\d[A-Z]{2}\b'
    postal_codes = re.findall(postal_pattern, text, re.IGNORECASE)
    if postal_codes:
        result['postal_code'] = postal_codes[0]
    amount_pattern = r'(\$[\d,]+\.?\d*)'
    amounts = re.findall(amount_pattern, text)
    if amounts:
        result['amount'] = amounts[0]
    return result

def is_valid_value(value: str) -> bool:
    if not value or len(value) < 2:
        return False
    if any(word in value.lower() for word in ['template', 'blank', 'form', '{{', '}}', 'placeholder']):
        return False
    if len(value) < 2 or len(value) > 200:
        return False
    return True

def match_form_fields_with_data(form_fields: List[str], client_data: Dict[str, str]) -> Dict[str, str]:
    matched_data = {}
    for field in form_fields:
        if field in client_data:
            matched_data[field] = client_data[field]
            continue
        matched_value = find_direct_match(field, client_data)
        if matched_value:
            matched_data[field] = matched_value
            continue
        matched_value = find_semantic_match(field, client_data)
        if matched_value:
            matched_data[field] = matched_value
            continue
        matched_value = find_field_variation(field, client_data)
        if matched_value:
            matched_data[field] = matched_value
            continue
        matched_data[field] = ""
    return matched_data

def find_direct_match(field_name: str, client_data: Dict[str, str]) -> str:
    field_lower = field_name.lower()
    for key, value in client_data.items():
        key_lower = key.lower()
        if field_lower in key_lower or key_lower in field_lower:
            return value
    return ""

def find_semantic_match(field_name: str, client_data: Dict[str, str]) -> str:
    field_lower = field_name.lower()
    best_match = ""
    best_score = 0.6
    for key, value in client_data.items():
        key_lower = key.lower()
        similarity = SequenceMatcher(None, field_lower, key_lower).ratio()
        if similarity > best_score:
            best_score = similarity
            best_match = value
    return best_match

def find_field_variation(field_name: str, client_data: Dict[str, str]) -> str:
    field_lower = field_name.lower()
    field_variations = {
        'name': ['name', 'full name', 'client name', 'applicant', 'person', 'individual'],
        'dob': ['dob', 'date of birth', 'birth', 'birth date', 'birthday', 'born'],
        'address': ['address', 'residence', 'home', 'street', 'location', 'place'],
        'employment': ['employment', 'job', 'occupation', 'position', 'work', 'profession', 'career'],
        'income': ['income', 'salary', 'earnings', 'pay', 'wage', 'revenue', 'compensation'],
        'travel': ['travel', 'visited', 'countries', 'destinations', 'trips', 'journeys'],
        'email': ['email', 'e-mail', 'mail', 'contact'],
        'phone': ['phone', 'telephone', 'mobile', 'cell', 'contact number'],
        'postal_code': ['postal code', 'zip code', 'postcode', 'zip'],
        'amount': ['amount', 'sum', 'total', 'value', 'cost', 'price']
    }
    for base_field, variations in field_variations.items():
        if field_lower in variations:
            for variation in variations:
                if variation in client_data:
                    return client_data[variation]
    for base_field, variations in field_variations.items():
        for variation in variations:
            if field_lower in variation or variation in field_lower:
                if base_field in client_data:
                    return client_data[base_field]
    return ""

# --- PDF Filler ---
def fill_pdf(input_file, output_pdf, data):
    file_extension = os.path.splitext(input_file)[1].lower()
    if file_extension == '.pdf':
        return fill_pdf_form(input_file, output_pdf, data)
    elif file_extension == '.docx':
        return fill_docx_template(input_file, output_pdf, data)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")

def fill_pdf_form(input_pdf, output_pdf, data):
    try:
        template = PdfReader(input_pdf)
        for page in template.pages:
            annotations = page.get("/Annots")
            if annotations:
                for annotation in annotations:
                    if annotation['/T']:
                        key = annotation['/T'][1:-1]  # PDF field name
                        print("PDF field:", key)
                        # Try direct match, lower-case match, etc.
                        value = data.get(key) or data.get(key.lower()) or ""
                        if value:
                            annotation.update(PdfDict(V=f"{value}"))
        PdfWriter().write(output_pdf, template)
        return True
    except Exception as e:
        print(f"Error filling PDF form: {e}")
        return False

def fill_docx_template(input_docx, output_pdf, data):
    try:
        doc = docx.Document(input_docx)
        replace_placeholders_in_document(doc, data)
        temp_docx = input_docx.replace('.docx', '_filled.docx')
        doc.save(temp_docx)
        convert_docx_to_pdf(temp_docx, output_pdf)
        if os.path.exists(temp_docx):
            os.remove(temp_docx)
        return True
    except Exception as e:
        print(f"Error filling DOCX template: {e}")
        return False

def replace_placeholders_in_document(doc, data):
    for paragraph in doc.paragraphs:
        paragraph.text = replace_placeholders_in_text(paragraph.text, data)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    paragraph.text = replace_placeholders_in_text(paragraph.text, data)

def replace_placeholders_in_text(text, data):
    pattern = r'\{\{([^}]+)\}\}'
    def replace_match(match):
        field_name = match.group(1).strip().lower()
        for key, value in data.items():
            if key.lower() == field_name:
                return str(value)
        for key, value in data.items():
            if field_name in key.lower() or key.lower() in field_name:
                return str(value)
        return match.group(0)
    return re.sub(pattern, replace_match, text)

def convert_docx_to_pdf(docx_path, pdf_path):
    try:
        doc = docx.Document(docx_path)
        text_content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        create_simple_pdf_from_text(text_content, pdf_path)
    except Exception as e:
        print(f"Error converting DOCX to PDF: {e}")
        create_simple_pdf_from_text(["Document conversion failed. Please check the original file."], pdf_path)

def create_simple_pdf_from_text(text_lines, pdf_path):
    try:
        doc = fitz.open()
        page = doc.new_page()
        y_position = 50
        for line in text_lines:
            if line.strip():
                page.insert_text((50, y_position), line, fontsize=12)
                y_position += 20
        doc.save(pdf_path)
        doc.close()
    except Exception as e:
        print(f"Error creating simple PDF: {e}")
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), "PDF generation completed. Check the filled form.", fontsize=12)
        doc.save(pdf_path)
        doc.close()

# --- AI Extractor (Gemini) ---
genai.configure(api_key="AIzaSyBRAKfYEnbImVZOEESX7KuIA8Op5mWI9js")  # Replace with your real key

def extract_required_info(raw_text):
    prompt = f"""
    Extract the following information from this text:
    - Name
    - Date of Birth
    - Address
    - Employment
    - Income
    - Travel history
    Provide the result in a JSON format with keys: Name, DOB, Address, Employment, Income, Travel.
    Text:
    {raw_text}
    """
    model = genai.GenerativeModel("gemini-2.5-flash-lite")
    response = model.generate_content(prompt)
    try:
        return json.loads(response.text)
    except:
        return {
            "Name": "",
            "DOB": "",
            "Address": "",
            "Employment": "",
            "Income": "",
            "Travel": ""
        }

# --- FastAPI Routes ---
# @app.route("/", methods=["GET"])
# def home():
#     return render_template("index.html")

@app.route("/process", methods=["POST"])
def process_form():
    global form_fields, extracted_data
    try:
        template_file = request.files['template_file']
        client_file = request.files['client_file']

        template_path = os.path.join(UPLOAD_DIR, template_file.filename)
        client_path = os.path.join(UPLOAD_DIR, client_file.filename)

        template_file.save(template_path)
        client_file.save(client_path)

        template_text = extract_text_from_file(template_path)
        form_fields = extract_form_fields_from_template(template_text)

        client_text = extract_text_from_file(client_path)
        client_data = extract_all_data_from_client_document(client_text)

        extracted_data = match_form_fields_with_data(form_fields, client_data)

        return jsonify({
            "template_path": template_path,
            "fields": extracted_data
        })
    except Exception as e:
        print(f"Error in /process: {e}")
        return jsonify({"error": f"Error processing files: {str(e)}"}), 500

@app.route("/generate_pdf", methods=["POST"])
def generate_pdf_handler():
    global form_fields
    try:
        template_path = request.form.get("template_path")
        data = {}
        for field in form_fields:
            field_value = request.form.get(field, "")
            data[field] = field_value  # ✅ use field name as-is

        print(f"Filling PDF with data: {data}")
        output_file = os.path.join(FILLED_DIR, "filled_form.pdf")

        success = fill_pdf(template_path, output_file, data)

        if success:
             return render_template("<h3>✅ PDF generated successfully.</h3><a href='/download-pdf'>⬇️ Download PDF</a>")
        else:
             return render_template("<h3>❌ PDF generation failed.</h3><p>Please check your template and try again.</p>")

    except Exception as e:
        print("❌ Error generating PDF:", e)
        import traceback
        traceback.print_exc()  # ✅ shows full error details in console
        return render_template("error.html", error_message=f"Error generating PDF: {str(e)}")


@app.route("/download-pdf", methods=["GET"])
def download_pdf():
    try:
        return send_file(os.path.join(FILLED_DIR, "filled_form.pdf"), as_attachment=True, download_name="filled_form.pdf")
    except Exception as e:
        print(f"Error sending PDF: {e}")
        return render_template("error.html", error_message="PDF file not found or could not be sent.")


# for cases 
@app.route("/upload-case", methods=["POST"])
def upload_case():
    if 'cases[]' not in request.files:
        return jsonify({"error": "No case files provided"}), 400

    uploaded_files = request.files.getlist('cases[]')
    if not uploaded_files:
        return jsonify({"error": "Empty case file list"}), 400

    responses = []

    for file in uploaded_files:
        if file.filename == '':
            responses.append({"error": "Empty filename"})
            continue

        try:
            filename = file.filename
            ext = os.path.splitext(filename)[1].lower()
            filetype = ext[1:]  # remove dot
            file_bytes = file.read()

            # Extract text
            if ext == ".txt":
                extracted_text = file_bytes.decode("utf-8")
            elif ext == ".docx":
                doc = Document(BytesIO(file_bytes))
                extracted_text = "\n".join([para.text for para in doc.paragraphs])
            elif ext == ".pdf":
                pdf = fitz.open(stream=file_bytes, filetype="pdf")
                extracted_text = "\n".join([page.get_text() for page in pdf])
            elif ext in [".jpg", ".jpeg", ".png"]:
                image = Image.open(BytesIO(file_bytes))
                extracted_text = pytesseract.image_to_string(image)
            else:
                responses.append({"filename": filename, "error": "Unsupported file type."})
                continue

            # You may want to extract more metadata for cases, e.g. title, description, etc.
            # For now, we'll use filename as title and extracted_text as description.
            title = filename
            description = extracted_text
            case_type = "Unknown"  # You can set this from frontend
            status = "Open"
            created_by = request.form.get("created_by", None)
            user_id = request.form.get("user_id", None)

            # Save to DB
            conn = get_connection()
            cur = conn.cursor()
            cur.execute("""
                INSERT INTO cases (title, description, case_type, status, created_by, user_id)
                VALUES (%s, %s, %s, %s, %s, %s)
                RETURNING id
            """, (title, description, case_type, status, created_by, user_id))
            case_id = cur.fetchone()[0]
            conn.commit()
            cur.close()
            conn.close()

            responses.append({
                "filename": filename,
                "id": case_id,
                "text": extracted_text
            })

        except Exception as e:
            responses.append({
                "filename": file.filename,
                "error": str(e)
            })

    return jsonify({"results": responses})




# ================== CASES ================== #
# ================== CASES ================== #
@app.route("/cases", methods=["GET"])
def get_cases():
    try:
        print("🔍 Getting cases from database...")
        conn = get_connection()
        cur = conn.cursor()
        
        # First check if the cases table exists
        cur.execute("""
            SELECT EXISTS (
                SELECT FROM information_schema.tables 
                WHERE table_name = 'cases'
            );
        """)
        table_exists = cur.fetchone()[0]
        
        if not table_exists:
            print("⚠️ cases table does not exist!")
            cur.close()
            conn.close()
            return jsonify({"error": "cases table does not exist"}), 500
        
        print("✅ cases table exists, fetching data...")
        
        # Check table structure
        cur.execute("""
            SELECT column_name, data_type 
            FROM information_schema.columns 
            WHERE table_name = 'cases'
            ORDER BY ordinal_position;
        """)
        columns = cur.fetchall()
        print(f"📋 cases table columns: {[col[0] for col in columns]}")
        
        cur.execute("""
            SELECT 
                caseid, casename, description, casetype, casestatus, casepriority,
                uploadfile, createddate, openeddate, closeddate, duedate,
                courtname, hearingdate, judgename, casecategory, jurisdiction,
                casefee, billingstatus, customerid, document_type, document_preview
            FROM cases
        """)
        rows = cur.fetchall()
        print(f"📊 Found {len(rows)} cases in database")
        
        # Debug: Print first few rows
        for i, row in enumerate(rows[:3]):
            print(f"  Case {i+1}: {row}")
        
        cur.close()
        conn.close()

        cases = []
        for row in rows:
            cases.append({
                "id": row[0],
                "name": row[1],
                "description": row[2],
                "case_type": row[3],  # Changed from "type" to "case_type" to match frontend
                "status": row[4],
                "priority": row[5],
                "upload_file": row[6],
                "created_date": row[7].isoformat() if row[7] else None,
                "opened_date": row[8].isoformat() if row[8] else None,
                "closed_date": row[9].isoformat() if row[9] else None,
                "due_date": row[10].isoformat() if row[10] else None,
                "court": row[11],
                "hearing_date": row[12].isoformat() if row[12] else None,
                "judge": row[13],
                "category": row[14],
                "jurisdiction": row[15],
                "fee": float(row[16]) if row[16] else None,
                "billing_status": row[17],
                "customer_id": row[18],
                "document_type": row[19],
                "document_preview": row[20]
            })
        
        print(f"✅ Successfully processed {len(cases)} cases")
        return jsonify(cases)
    except Exception as e:
        print(f"❌ Error getting cases: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


@app.route("/cases", methods=["POST"])
def add_case():
    data = request.get_json()
    print(f"Received case data: {data}")  # Debug print
    try:
        conn = get_connection()
        cur = conn.cursor()
        
        # Check if customer exists
        customer_id = data.get("customer_id")
        if customer_id:
            cur.execute("SELECT CustomerID FROM Customer WHERE CustomerID = %s", (customer_id,))
            customer_exists = cur.fetchone()
            if not customer_exists:
                return jsonify({"error": f"Customer with ID {customer_id} does not exist. Please create a customer first."}), 400
        cur.execute("""
            INSERT INTO Cases 
            (CaseName, Description, CaseType, CaseStatus, CasePriority, UploadFile,
             OpenedDate, ClosedDate, DueDate, CourtName, HearingDate, JudgeName,
             CaseCategory, Jurisdiction, CaseFee, BillingStatus, CustomerID, document_type, document_preview)
            VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s)
            RETURNING CaseID
        """, (
            data.get("name"), data.get("description"), data.get("type"), 
            data.get("status", "Open"), data.get("priority", "Medium"), data.get("upload_file"),
            data.get("opened_date"), data.get("closed_date"), data.get("due_date"),
            data.get("court"), data.get("hearing_date"), data.get("judge"),
            data.get("category"), data.get("jurisdiction"), data.get("fee"),
            data.get("billing_status", "Pending"), data.get("customer_id"), data.get("document_type", "Case"),
            data.get("document_preview")
        ))
        case_id = cur.fetchone()[0]
        conn.commit(); cur.close(); conn.close()
        return jsonify({"id": case_id, "message": "Case added successfully"})
    except Exception as e:
        print(f"Error adding case: {e}")  # Add debugging
        return jsonify({"error": str(e)}), 500


@app.route("/upload-case-files", methods=["POST"])
def upload_case_files():
    """Handle file uploads for cases and save to Documents table with URLs"""
    print(f"Upload case files called with case_id: {request.form.get('case_id')}")
    print(f"Files in request: {list(request.files.keys())}")
    
    if 'files[]' not in request.files:
        return jsonify({"error": "No files provided"}), 400

    uploaded_files = request.files.getlist('files[]')
    case_id = request.form.get('case_id')
    
    print(f"Number of files: {len(uploaded_files)}")
    
    if not uploaded_files:
        return jsonify({"error": "Empty file list"}), 400

    responses = []
    for file in uploaded_files:
        if file.filename == '':
            responses.append({"error": "Empty filename"})
            continue

        try:
            # Save file and get URL
            file_url = save_file_and_get_url(file, 'cases')
            if not file_url:
                responses.append({"filename": file.filename, "error": "Failed to save file"})
                continue

            # Extract text from the file
            file_bytes = file.read()
            file.seek(0)  # Reset file pointer for saving
            extracted_text = ""
            
            try:
                ext = os.path.splitext(file.filename)[1].lower()
                if ext == ".txt":
                    extracted_text = file_bytes.decode("utf-8")
                elif ext == ".docx":
                    doc = Document(BytesIO(file_bytes))
                    extracted_text = "\n".join([para.text for para in doc.paragraphs])
                elif ext == ".pdf":
                    pdf = fitz.open(stream=file_bytes, filetype="pdf")
                    extracted_text = "\n".join([page.get_text() for page in pdf])
                elif ext in [".jpg", ".jpeg", ".png"]:
                    image = Image.open(BytesIO(file_bytes))
                    extracted_text = pytesseract.image_to_string(image)
                elif ext in [".ppt", ".pptx"]:
                    try:
                        # Try to import pptx, but don't fail if not available
                        try:
                            from pptx import Presentation  # type: ignore
                            prs = Presentation(BytesIO(file_bytes))
                            extracted_text = []
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        extracted_text.append(shape.text)
                            extracted_text = "\n".join(extracted_text)
                        except ImportError:
                            extracted_text = "PowerPoint file - text extraction requires python-pptx package. Install with: pip install python-pptx"
                    except Exception as e:
                        extracted_text = f"PowerPoint file - extraction error: {str(e)}"
            except Exception as e:
                print(f"Warning: Could not extract text from {file.filename}: {e}")
                extracted_text = ""

            # Save to Documents table with case DocumentType and extracted text
            conn = get_connection(); cur = conn.cursor()
            # Ensure preview column exists
            try:
                cur.execute("ALTER TABLE Documents ADD COLUMN IF NOT EXISTS document_preview TEXT;")
            except Exception:
                pass
            try:
                cur.execute(
                    """
                    INSERT INTO Documents (FileName, FileURL, DocumentType, CaseID, extracted_text, document_preview, content, filetype)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
                    RETURNING DocumentID
                    """,
                    (file.filename, file_url, 'Case', case_id, extracted_text, extracted_text, psycopg2.Binary(file_bytes), os.path.splitext(file.filename)[1][1:])
                )
            except Exception:
                # Fallback if column not present
                cur.execute(
                    """
                INSERT INTO Documents (FileName, FileURL, DocumentType, CaseID, extracted_text, content, filetype)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
                RETURNING DocumentID
                    """,
                    (file.filename, file_url, 'Case', case_id, extracted_text, psycopg2.Binary(file_bytes), os.path.splitext(file.filename)[1][1:])
                )
            doc_id = cur.fetchone()[0]
            conn.commit(); cur.close(); conn.close()

            responses.append({"filename": file.filename, "id": doc_id, "url": file_url, "text": extracted_text})

        except Exception as e:
            responses.append({"filename": file.filename, "error": str(e)})

    return jsonify({"results": responses})


@app.route("/upload-customer-files", methods=["POST"])
def upload_customer_files():
    """Handle file uploads for customers and save to Documents table with URLs"""
    print(f"Upload customer files called with customer_id: {request.form.get('customer_id')}")
    print(f"Files in request: {list(request.files.keys())}")
    
    if 'files[]' not in request.files:
        return jsonify({"error": "No files provided"}), 400

    uploaded_files = request.files.getlist('files[]')
    customer_id = request.form.get('customer_id')
    
    print(f"Number of files: {len(uploaded_files)}")
    
    if not uploaded_files:
        return jsonify({"error": "Empty file list"}), 400

    responses = []
    for file in uploaded_files:
        if file.filename == '':
            responses.append({"error": "Empty filename"})
            continue

        try:
            # Save file and get URL
            file_url = save_file_and_get_url(file, 'customers')
            if not file_url:
                responses.append({"filename": file.filename, "error": "Failed to save file"})
                continue

            # Extract text from the file
            file_bytes = file.read()
            file.seek(0)  # Reset file pointer for saving
            extracted_text = ""
            
            try:
                ext = os.path.splitext(file.filename)[1].lower()
                if ext == ".txt":
                    extracted_text = file_bytes.decode("utf-8")
                elif ext == ".docx":
                    doc = Document(BytesIO(file_bytes))
                    extracted_text = "\n".join([para.text for para in doc.paragraphs])
                elif ext == ".pdf":
                    pdf = fitz.open(stream=file_bytes, filetype="pdf")
                    extracted_text = "\n".join([page.get_text() for page in pdf])
                elif ext in [".jpg", ".jpeg", ".png"]:
                    image = Image.open(BytesIO(file_bytes))
                    extracted_text = pytesseract.image_to_string(image)
            except Exception as e:
                print(f"Warning: Could not extract text from {file.filename}: {e}")
                extracted_text = ""

            # Save to Documents table with customer DocumentType and extracted text
            conn = get_connection(); cur = conn.cursor()
            try:
                cur.execute("ALTER TABLE Documents ADD COLUMN IF NOT EXISTS document_preview TEXT;")
            except Exception:
                pass
            try:
                cur.execute(
                    """
                    INSERT INTO Documents (FileName, FileURL, DocumentType, CustomerID, extracted_text, document_preview, content, filetype)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
                    RETURNING DocumentID
                    """,
                    (file.filename, file_url, 'Customer', customer_id, extracted_text, extracted_text, psycopg2.Binary(file_bytes), os.path.splitext(file.filename)[1][1:])
                )
            except Exception:
                cur.execute(
                    """
                INSERT INTO Documents (FileName, FileURL, DocumentType, CustomerID, extracted_text, content, filetype)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
                RETURNING DocumentID
                    """,
                    (file.filename, file_url, 'Customer', customer_id, extracted_text, psycopg2.Binary(file_bytes), os.path.splitext(file.filename)[1][1:])
                )
            doc_id = cur.fetchone()[0]
            conn.commit(); cur.close(); conn.close()

            responses.append({"filename": file.filename, "id": doc_id, "url": file_url, "text": extracted_text})

        except Exception as e:
            responses.append({"filename": file.filename, "error": str(e)})

    return jsonify({"results": responses})


# ================== CUSTOMERS ================== #
@app.route("/customers", methods=["GET"])
def get_customers():
    try:
        print("🔍 Getting customers from database...")
        conn = get_connection()
        cur = conn.cursor()
        
        # First check if the customer table exists
        cur.execute("""
            SELECT EXISTS (
                SELECT FROM information_schema.tables 
                WHERE table_name = 'customer'
            );
        """)
        table_exists = cur.fetchone()[0]
        
        if not table_exists:
            print("⚠️ customer table does not exist!")
            cur.close()
            conn.close()
            return jsonify({"error": "customer table does not exist"}), 500
        
        print("✅ customer table exists, fetching data...")
        
        # Check table structure
        cur.execute("""
            SELECT column_name, data_type 
            FROM information_schema.columns 
            WHERE table_name = 'customer'
            ORDER BY ordinal_position;
        """)
        columns = cur.fetchall()
        print(f"📋 customer table columns: {[col[0] for col in columns]}")
        
        cur.execute("""
            SELECT customerid, name, contactinfo, phone, email, address, profilefile,
                   dateofbirth, gender, cnic, occupation, companyname,
                   emergencycontactname, emergencycontactphone, notes,
                   createddate, updateddate, document_type
            FROM customer
        """)
        rows = cur.fetchall()
        print(f"📊 Found {len(rows)} customers in database")
        
        # Debug: Print first few rows
        for i, row in enumerate(rows[:3]):
            print(f"  Customer {i+1}: {row}")
        
        cur.close()
        conn.close()

        customers = []
        for row in rows:
            customers.append({
                "id": row[0],
                "name": row[1],
                "contact_info": row[2],
                "phone": row[3],
                "email": row[4],
                "address": row[5],
                "profile_file": row[6],
                "date_of_birth": row[7].isoformat() if row[7] else None,
                "gender": row[8],
                "cnic": row[9],
                "occupation": row[10],
                "company_name": row[11],
                "emergency_contact_name": row[12],
                "emergency_contact_phone": row[13],
                "notes": row[14],
                "created_date": row[15].isoformat() if row[15] else None,
                "updated_date": row[16].isoformat() if row[16] else None,
                "document_type": row[17]
            })
        
        print(f"✅ Successfully processed {len(customers)} customers")
        return jsonify(customers)
    except Exception as e:
        print(f"❌ Error getting customers: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


@app.route("/customers", methods=["POST"])
def add_customer():
    data = request.get_json()
    try:
        conn = get_connection()
        cur = conn.cursor()
        cur.execute("""
            INSERT INTO customer 
            (name, contactinfo, phone, email, address, profilefile, dateofbirth, 
             gender, cnic, occupation, companyname, emergencycontactname, 
             emergencycontactphone, notes, document_type)
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
            RETURNING customerid
        """, (
            data.get("name"), data.get("contact_info"), data.get("phone"), 
            data.get("email"), data.get("address"), data.get("profile_file"),
            data.get("date_of_birth"), data.get("gender"), data.get("cnic"), 
            data.get("occupation"), data.get("company_name"), 
            data.get("emergency_contact_name"), data.get("emergency_contact_phone"),
            data.get("notes"), data.get("document_type", "Customer")
        ))
        customer_id = cur.fetchone()[0]
        conn.commit(); cur.close(); conn.close()
        return jsonify({"id": customer_id, "message": "Customer added successfully"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/customers/<int:customer_id>", methods=["DELETE"])
def delete_customer(customer_id):
    """Delete a customer by ID"""
    try:
        conn = get_connection()
        cur = conn.cursor()
        
        # Check if customer exists
        cur.execute("SELECT customerid FROM customer WHERE customerid = %s", (customer_id,))
        if not cur.fetchone():
            cur.close(); conn.close()
            return jsonify({"error": "Customer not found"}), 404
        
        # Delete customer
        cur.execute("DELETE FROM customer WHERE customerid = %s", (customer_id,))
        conn.commit()
        cur.close(); conn.close()
        
        return jsonify({"message": "Customer deleted successfully"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/customers/<int:customer_id>", methods=["PUT"])
def update_customer(customer_id):
    """Update a customer by ID"""
    data = request.get_json()
    try:
        conn = get_connection()
        cur = conn.cursor()
        
        # Check if customer exists
        cur.execute("SELECT customerid FROM customer WHERE customerid = %s", (customer_id,))
        if not cur.fetchone():
            cur.close(); conn.close()
            return jsonify({"error": "Customer not found"}), 404
        
        # Update customer
        cur.execute("""
            UPDATE customer SET 
                name = %s, contactinfo = %s, phone = %s, email = %s, address = %s,
                profilefile = %s, dateofbirth = %s, gender = %s, cnic = %s,
                occupation = %s, companyname = %s, emergencycontactname = %s,
                emergencycontactphone = %s, notes = %s, document_type = %s
            WHERE customerid = %s
        """, (
            data.get("name"), data.get("contact_info"), data.get("phone"), 
            data.get("email"), data.get("address"), data.get("profile_file"),
            data.get("date_of_birth"), data.get("gender"), data.get("cnic"), 
            data.get("occupation"), data.get("company_name"), 
            data.get("emergency_contact_name"), data.get("emergency_contact_phone"),
            data.get("notes"), data.get("document_type", "Customer"), customer_id
        ))
        
        conn.commit()
        cur.close(); conn.close()
        return jsonify({"message": "Customer updated successfully"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/cases/<int:case_id>", methods=["DELETE"])
def delete_case(case_id):
    """Delete a case by ID"""
    try:
        conn = get_connection()
        cur = conn.cursor()
        
        # Check if case exists
        cur.execute("SELECT caseid FROM cases WHERE caseid = %s", (case_id,))
        if not cur.fetchone():
            cur.close(); conn.close()
            return jsonify({"error": "Case not found"}), 404
        
        # Delete case
        cur.execute("DELETE FROM cases WHERE caseid = %s", (case_id,))
        conn.commit()
        cur.close(); conn.close()
        
        return jsonify({"message": "Case deleted successfully"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/cases/<int:case_id>", methods=["PUT"])
def update_case(case_id):
    """Update a case by ID"""
    data = request.get_json()
    try:
        conn = get_connection()
        cur = conn.cursor()
        
        # Check if case exists
        cur.execute("SELECT CaseID FROM Cases WHERE CaseID = %s", (case_id,))
        if not cur.fetchone():
            cur.close(); conn.close()
            return jsonify({"error": "Case not found"}), 404
        
        # Check if customer exists
        customer_id = data.get("customer_id")
        if customer_id:
            cur.execute("SELECT CustomerID FROM Customer WHERE CustomerID = %s", (customer_id,))
            customer_exists = cur.fetchone()
            if not customer_exists:
                cur.close(); conn.close()
                return jsonify({"error": f"Customer with ID {customer_id} does not exist. Please create a customer first."}), 400
        
        # Update case
        cur.execute("""
            UPDATE Cases SET 
                CaseName = %s, Description = %s, CaseType = %s, CaseStatus = %s,
                CasePriority = %s, UploadFile = %s, OpenedDate = %s, ClosedDate = %s,
                DueDate = %s, CourtName = %s, HearingDate = %s, JudgeName = %s,
                CaseCategory = %s, Jurisdiction = %s, CaseFee = %s, BillingStatus = %s,
                CustomerID = %s, document_type = %s, document_preview = %s
            WHERE CaseID = %s
        """, (
            data.get("name"), data.get("description"), data.get("type"), 
            data.get("status", "Open"), data.get("priority", "Medium"), data.get("upload_file"),
            data.get("opened_date"), data.get("closed_date"), data.get("due_date"),
            data.get("court"), data.get("hearing_date"), data.get("judge"),
            data.get("category"), data.get("jurisdiction"), data.get("fee"),
            data.get("billing_status", "Pending"), data.get("customer_id"), 
            data.get("document_type", "Case"), data.get("document_preview"), case_id
        ))
        
        conn.commit()
        cur.close(); conn.close()
        return jsonify({"message": "Case updated successfully"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/")
def home():
    return render_template("index.html")

@app.route("/create-test-data", methods=["POST"])
def create_test_data():
    """Create some test data for debugging"""
    try:
        conn = get_connection()
        cur = conn.cursor()
        
        # Create test customer
        cur.execute("""
            INSERT INTO Customer (Name, Phone, Email, Address, document_type)
            VALUES (%s, %s, %s, %s, %s)
            RETURNING CustomerID
        """, ("Test Customer", "1234567890", "test@example.com", "123 Test Street", "Customer"))
        customer_id = cur.fetchone()[0]
        
        # Create test case
        cur.execute("""
            INSERT INTO Cases (CaseName, Description, CaseType, CaseStatus, CasePriority, CustomerID, document_type)
            VALUES (%s, %s, %s, %s, %s, %s, %s)
            RETURNING CaseID
        """, ("Test Case", "This is a test case for debugging", "Civil", "Open", "Medium", customer_id, "Case"))
        
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "message": "Test data created successfully",
            "customer_id": customer_id
        })
    except Exception as e:
        return jsonify({
            "success": False,
            "message": f"Error creating test data: {str(e)}"
        }), 500

@app.route("/test")
def test():
    """Simple test route to check if backend is working"""
    try:
        conn = get_connection()
        cur = conn.cursor()
        
        # Check if tables exist
        cur.execute("""
            SELECT table_name 
            FROM information_schema.tables 
            WHERE table_name IN ('Cases', 'Customer', 'Documents')
        """)
        existing_tables = [row[0] for row in cur.fetchall()]
        
        # Get row counts
        case_count = 0
        customer_count = 0
        document_count = 0
        
        if 'Cases' in existing_tables:
            cur.execute("SELECT COUNT(*) FROM Cases")
            case_count = cur.fetchone()[0]
        
        if 'Customer' in existing_tables:
            cur.execute("SELECT COUNT(*) FROM Customer")
            customer_count = cur.fetchone()[0]
            
        if 'Documents' in existing_tables:
            cur.execute("SELECT COUNT(*) FROM Documents")
            document_count = cur.fetchone()[0]
        
        cur.close()
        conn.close()
        
        return jsonify({
            "status": "success",
            "message": "Backend is working",
            "existing_tables": existing_tables,
            "case_count": case_count,
            "customer_count": customer_count,
            "document_count": document_count
        })
    except Exception as e:
        return jsonify({
            "status": "error",
            "message": f"Database connection failed: {str(e)}"
        }), 500

@app.route("/uploads/<path:filename>")
def uploaded_file(filename):
    """Serve uploaded files"""
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

# ================== DOCUMENTS ================== #
@app.route("/documents", methods=["GET"])
def get_documents():
    """Get all documents from the Documents table"""
    try:
        conn = get_connection()
        cur = conn.cursor()
        cur.execute("""
            SELECT DocumentID, FileName, FileURL, DocumentType, 
                   CustomerID, CaseID, UploadDate
            FROM Documents
            ORDER BY UploadDate DESC
        """)
        rows = cur.fetchall()
        cur.close(); conn.close()

        documents = []
        for row in rows:
            documents.append({
                "document_id": row[0],
                "filename": row[1],
                "file_url": row[2],
                "document_type": row[3],
                "customer_id": row[4],
                "case_id": row[5],
                "upload_date": row[6].isoformat() if row[6] else None
            })

        return jsonify(documents)
    except Exception as e:
        print(f"Error getting documents: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/documents/<int:document_id>", methods=["DELETE"])
def delete_document(document_id):
    """Delete a document by ID"""
    try:
        conn = get_connection()
        cur = conn.cursor()
        
        # Get file info before deletion
        cur.execute("SELECT FileURL FROM Documents WHERE DocumentID = %s", (document_id,))
        result = cur.fetchone()
        
        if not result:
            cur.close(); conn.close()
            return jsonify({"error": "Document not found"}), 404
        
        file_url = result[0]
        
        # Delete from database
        cur.execute("DELETE FROM Documents WHERE DocumentID = %s", (document_id,))
        conn.commit()
        cur.close(); conn.close()
        
        # Delete physical file if it exists
        if file_url and file_url.startswith('/uploads/'):
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], file_url.replace('/uploads/', ''))
            if os.path.exists(file_path):
                os.remove(file_path)
        
        return jsonify({"message": "Document deleted successfully"})
    except Exception as e:
        print(f"Error deleting document: {e}")
        return jsonify({"error": str(e)}), 500




if __name__ == "__main__":
    # Ensure all required tables exist
    ensure_cases_customers_tables()
    ensure_documents_table()
    ensure_signup_tables()
    
    Timer(1, lambda: webbrowser.open("http://127.0.0.1:5050")).start()
    app.run(debug=True, port=5050)

    # port = int(os.environ.get("PORT", 5000))  # Azure PORT set karega
    # app.run(host="0.0.0.0", port=port)


